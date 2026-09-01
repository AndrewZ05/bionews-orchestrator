"""
Generate Identity Hub PowerPoint Presentation
Explains the entire identity resolution process in plain language.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# 16:9 aspect ratio
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Colors
BLUE = RGBColor(0x1a, 0x73, 0xe8)
GREEN = RGBColor(0x34, 0xa8, 0x53)
RED = RGBColor(0xea, 0x43, 0x35)
ORANGE = RGBColor(0xfb, 0xbc, 0x04)
PURPLE = RGBColor(0x9c, 0x27, 0xb0)
GRAY = RGBColor(0x60, 0x7d, 0x8b)
DARK = RGBColor(0x33, 0x33, 0x33)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_BLUE = RGBColor(0xe8, 0xf0, 0xfe)
LIGHT_GREEN = RGBColor(0xe6, 0xf4, 0xea)
LIGHT_RED = RGBColor(0xfc, 0xe8, 0xe6)


def add_title_slide(title, subtitle=""):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    # Background
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = BLUE
    bg.line.fill.background()
    # Title
    box = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(11.333), Inches(1.5))
    p = box.text_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(48)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    if subtitle:
        box2 = slide.shapes.add_textbox(Inches(1), Inches(4.2), Inches(11.333), Inches(0.75))
        p2 = box2.text_frame.paragraphs[0]
        p2.text = subtitle
        p2.font.size = Pt(24)
        p2.font.color.rgb = RGBColor(0xBB, 0xDE, 0xFB)
        p2.alignment = PP_ALIGN.CENTER


def add_section_slide(title):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = RGBColor(0xF5, 0xF5, 0xF5)
    bg.line.fill.background()
    box = slide.shapes.add_textbox(Inches(1), Inches(3), Inches(11.333), Inches(1.5))
    p = box.text_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = BLUE
    p.alignment = PP_ALIGN.CENTER


def add_slide(title, bullets, note=""):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    # Title bar
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.2))
    bar.fill.solid()
    bar.fill.fore_color.rgb = BLUE
    bar.line.fill.background()
    tbox = slide.shapes.add_textbox(Inches(0.5), Inches(0.15), Inches(12), Inches(0.9))
    tp = tbox.text_frame.paragraphs[0]
    tp.text = title
    tp.font.size = Pt(32)
    tp.font.bold = True
    tp.font.color.rgb = WHITE
    # Bullets
    content = slide.shapes.add_textbox(Inches(0.8), Inches(1.6), Inches(11.5), Inches(5.2))
    tf = content.text_frame
    tf.word_wrap = True
    for i, bullet in enumerate(bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = bullet
        p.font.size = Pt(20)
        p.font.color.rgb = DARK
        p.space_after = Pt(12)
        p.level = 0
        if bullet.startswith("  "):
            p.level = 1
            p.font.size = Pt(18)
            p.font.color.rgb = GRAY
    if note:
        nbox = slide.shapes.add_textbox(Inches(0.8), Inches(6.8), Inches(11.5), Inches(0.5))
        np = nbox.text_frame.paragraphs[0]
        np.text = note
        np.font.size = Pt(14)
        np.font.color.rgb = GRAY
        np.font.italic = True


def add_table_slide(title, headers, rows, note=""):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.2))
    bar.fill.solid()
    bar.fill.fore_color.rgb = BLUE
    bar.line.fill.background()
    tbox = slide.shapes.add_textbox(Inches(0.5), Inches(0.15), Inches(12), Inches(0.9))
    tp = tbox.text_frame.paragraphs[0]
    tp.text = title
    tp.font.size = Pt(32)
    tp.font.bold = True
    tp.font.color.rgb = WHITE

    cols = len(headers)
    total_rows = len(rows) + 1
    table_width = Inches(12)
    table = slide.shapes.add_table(total_rows, cols, Inches(0.667), Inches(1.5), table_width, Inches(0.4 * total_rows)).table

    # Style header
    for j, h in enumerate(headers):
        cell = table.cell(0, j)
        cell.text = h
        cell.fill.solid()
        cell.fill.fore_color.rgb = BLUE
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(16)
            p.font.bold = True
            p.font.color.rgb = WHITE

    # Style rows
    for i, row in enumerate(rows):
        for j, val in enumerate(row):
            cell = table.cell(i + 1, j)
            cell.text = str(val)
            cell.fill.solid()
            cell.fill.fore_color.rgb = WHITE if i % 2 == 0 else LIGHT_BLUE
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(14)
                p.font.color.rgb = DARK

    if note:
        nbox = slide.shapes.add_textbox(Inches(0.8), Inches(6.8), Inches(11.5), Inches(0.5))
        np = nbox.text_frame.paragraphs[0]
        np.text = note
        np.font.size = Pt(14)
        np.font.color.rgb = GRAY
        np.font.italic = True


# ============================================================================
# SLIDES
# ============================================================================

# Slide 1: Title
add_title_slide(
    "BioNews Identity Hub",
    "Cross-Platform Identity Resolution Platform  |  April 2026 (v3.5)"
)

# Slide 2: What Is It?
add_slide("What Is the Identity Hub?", [
    "A system that connects all the different ways we see the same person",
    "When someone reads an email, visits our website, and fills out a survey...",
    "  ...the Identity Hub knows it's all the SAME person",
    "",
    "It produces a single persistent ID (bn_id) for each person",
    "  Example: BN_Xy7a2b3c4d5e6f = one person across all touchpoints",
    "",
    "This enables: unified profiles, personalized content, audience targeting, analytics",
])

# Slide 3: Scale
add_table_slide("Current Scale (April 2026)", [
    "Metric", "Value", "Description"
], [
    ["Total Persons", "3,518,018", "Unique people identified"],
    ["Known Users (Tier 1)", "770,458", "Have email, NPI, mc_euid, or bionews_uk"],
    ["Anonymous Visitors (Tier 2)", "2,747,560", "Have bnfpvid only (browser ID)"],
    ["Identifier Linkages (xref rows)", "18.5M", "Identifier -> bn_id rows"],
    ["Identity Edges (hub rows)", "39.2M", "Connections between identifiers"],
    ["HCP Clusters", "340,643", "bn_ids flagged is_hcp=TRUE"],
    ["LimeSurvey Respondents", "4,142", "Survey participants linked to profiles"],
    ["Identifier Types", "17", "Different types of IDs tracked"],
])

# Slide 4: How It Works - Overview
add_slide("How Does It Work?", [
    "Step 1: COLLECT  --  Gather identifiers from every data source",
    "  Mailchimp emails, WordPress logins, GA4 browser IDs, NPI numbers, cookies...",
    "",
    "Step 2: CONNECT  --  Find links between identifiers",
    "  Same email in Mailchimp AND WordPress? Same person.",
    "  Same browser visited site AND clicked Mailchimp email? Same person.",
    "",
    "Step 3: STITCH  --  Merge all connected identifiers into one cluster",
    "  Union-Find algorithm groups all linked IDs into one bn_id",
    "",
    "Step 4: PUBLISH  --  Write the resolved identities to BigQuery tables",
    "  bn_id_xref: look up any identifier -> get the person's bn_id",
])

# Slide 5: Data Sources
add_section_slide("Where Do the Identifiers Come From?")

# Slide 6: Source Table
add_table_slide("Data Sources & Connectors", [
    "Source", "What It Provides", "Edge Count", "Type"
], [
    ["BN_Acceptor (acceptor)", "Browser ID, cookies, device, AIM payload", "34.9M", "Browser"],
    ["BN_Acceptor (bio_acceptor)", "Preprocessed identity columns", "1.9M", "Browser"],
    ["Mailchimp", "Email, subscriber hash, mc_euid, phone", "1.0M", "Reference"],
    ["GA4 Events", "Session co-occurrence + URL mc_euid", "881K", "Browser"],
    ["NPI Registry", "NPI number + email/phone bridges", "312K", "Reference"],
    ["DMD Audiences", "Verified HCP email/mc_euid -> NPI", "93K", "Reference"],
    ["WordPress Users", "wp_user_id + email bridge", "54K", "Reference"],
    ["AIM Clickstream", "aim_dgid -> NPI (weekly webfeed)", "27K", "Reference"],
    ["LimeSurvey", "Survey respondent email + session bridge", "5.2K", "Reference"],
    ["Gravity Forms", "Form email -> GA4 session via IP+time", "4.0K", "Event"],
], note="Counts are total rows in bn_id_hub by source_system as of April 2026.")

# Slide 7: The Stitching Chain
add_slide("The Identity Chain", [
    "email  <-->  subscriber_hash  (Mailchimp bridge, deterministic)",
    "  |",
    "email  <-->  mc_euid  (Mailchimp encrypted ID, deterministic)",
    "  |",
    "mc_euid  <-->  bnfpvid  (localStorage co-occurrence on website)",
    "  |",
    "bnfpvid  <-->  client_id  (GA4 client ID, same browser session)",
    "  |",
    "bnfpvid  <-->  aim_dgid  (AIM healthcare professional signal)",
    "  |",
    "aim_dgid  <-->  npi_number  (AIM Clickstream, verified HCP)",
    "",
    "Result: email -> mc_euid -> bnfpvid -> client_id -> aim_dgid -> npi_number",
    "One person, six different identifiers, all connected.",
])

# Slide 8: Quality Controls
add_section_slide("Quality & Safety Controls")

# Slide 9: Quality Protections
add_slide("How We Prevent Bad Merges", [
    "1. FANOUT THRESHOLDS  --  If one identifier connects to too many others, it's suspicious",
    "  Email in 10+ clusters? Probably a shared/invalid address. Excluded.",
    "",
    "2. SHARED WORKSTATION DETECTION  --  Clinic computers where multiple HCPs log in",
    "  If one browser sees 3+ Mailchimp users, those bnfpvid edges are removed pre-stitch",
    "  If it sees 2+ users, the cluster gets flagged (is_shared_workstation = TRUE)",
    "",
    "3. CONFLICT GATE  --  Catches bridges connecting too many anchors",
    "  One bnfpvid connecting to 3+ different emails? Capped below stitch threshold.",
    "",
    "4. CONFIDENCE SCORING  --  Every edge has a confidence score (0-1)",
    "  Only edges >= 0.80 confidence participate in stitching",
    "  Probabilistic edges (IP matching, device fingerprint) stay below threshold unless repeated",
    "",
    "5. CLUSTER HEALTH SCORING  --  Every cluster gets a health score (0-100)",
    "  Multi-email, multi-NPI, excessive devices all reduce the score",
    "  is_suspicious = TRUE when score < 60",
])

# Slide 10: Confidence Model
add_table_slide("Confidence Model", [
    "Match Rule", "Base Confidence", "Can It Stitch?", "Description"
], [
    ["EMAIL_EXACT", "1.0", "Yes", "Deterministic email bridge"],
    ["LOCALSTORAGE_COOCCURRENCE", "1.0", "Yes", "Same browser, same page load"],
    ["MC_EUID_BRIDGE", "1.0", "Yes", "Mailchimp email-to-euid"],
    ["UTM_CROSSWALK", "0.90", "Yes", "mc_euid in URL matched to GA4 session"],
    ["MC_CLICK_IP_TIME", "0.85", "Yes", "IP + 5-min window matching"],
    ["COMPOUND_DEVICE_IP_TIME", "0.75", "With 2+ obs", "stat_id + IP + time (NEW)"],
    ["IP_DEVICE_TIME", "0.50", "No", "IP + time only, too noisy"],
    ["DEVICE_STAT_ID", "0.35", "No", "Browser fingerprint, audit only"],
], note="Confidence >= 0.80 required for stitching. Log2 boost for repeated observations.")

# Slide 11: Edge Aging
add_slide("Edge Aging & Decay", [
    "Not all connections stay strong forever.",
    "",
    "ANCHOR-AWARE DECAY: If one side of an edge is a known anchor (email, NPI, mc_euid),",
    "  the edge does NOT decay. The identity link is established.",
    "",
    "BROWSER-TO-BROWSER DECAY: Pure device-level connections weaken over time:",
    "  0-30 days: full strength (weight 1.0)",
    "  30-90 days: slight decay (weight 0.8)",
    "  90-180 days: moderate decay (weight 0.6)",
    "  180-365 days: significant decay (weight 0.4)",
    "  365+ days: expired (excluded from stitching)",
    "",
    "WHY: A user who visited 6 months ago with a different browser shouldn't",
    "  permanently merge with today's visitor based on an old cookie.",
])

# Slide 12: Output Tables
add_section_slide("Output Tables")

# Slide 13: Tables
add_table_slide("Identity Hub Output Tables", [
    "Table", "Rows", "Purpose", "How to Use"
], [
    ["bn_id_xref", "18.5M", "Look up any identifier -> bn_id", "Primary lookup surface"],
    ["bn_id_hub", "39.2M", "All edges with confidence", "Audit trail, debugging"],
    ["bn_id_neighbors", "54.7M", "Cluster member pairs", "Cluster inspection"],
    ["bn_id_node_index", "varies", "All stitchable identifiers", "Incremental subset discovery"],
    ["bn_id_identity_changes", "8.0K", "Merge/split audit log", "Ops monitoring"],
    ["bn_id_persistence", "4.7K", "bn_id redirect table", "Downstream bn_id stability"],
    ["bn_id_metrics", "1.0K", "Run statistics", "Health monitoring"],
    ["bn_id_manifest", "1", "Active run tracking", "Mixed-generation detection"],
])

# Slide 14: How to Query
add_slide("How to Look Up a Person", [
    "-- Find all identifiers for a person by email:",
    "SELECT * FROM identity_hub_data.bn_id_xref",
    "WHERE bn_id = (",
    "  SELECT bn_id FROM identity_hub_data.bn_id_xref",
    "  WHERE identifier_type = 'email'",
    "    AND identifier_value = 'user@example.com'",
    ")",
    "",
    "-- Find all neighbors (other identifiers in the same cluster):",
    "SELECT * FROM identity_hub_data.bn_id_neighbors",
    "WHERE bn_id = 'BN_Xy7a2b3c4d5e6f'",
    "",
    "-- Look up an NPI to find the person's email and browser IDs:",
    "SELECT identifier_type, identifier_value",
    "FROM identity_hub_data.bn_id_xref",
    "WHERE bn_id = (",
    "  SELECT bn_id FROM identity_hub_data.bn_id_xref",
    "  WHERE identifier_type = 'npi_number' AND identifier_value = '1234567890'",
    ")",
])

# Slide 15: Operations
add_section_slide("Operations & Monitoring")

# Slide 16: Running the Pipeline
add_slide("Running the Identity Hub", [
    "FULL REBUILD (monthly or after scoring changes):",
    "  python orchestrate.py --source identity_hub --env prod --refresh full",
    "  ~4.5 hours, rewrites all tables from scratch",
    "",
    "DAILY INCREMENTAL:",
    "  python orchestrate.py --source identity_hub --env prod --lookback 7",
    "  ~25-45 minutes, subset Union-Find on touched clusters",
    "",
    "WHEN TO REBUILD:",
    "  - Changed confidence thresholds or fanout limits",
    "  - Changed edge aging weights or identity ratios",
    "  - Added a new connector that needs retroactive edges",
    "",
    "WHEN INCREMENTAL IS FINE:",
    "  - New data arriving daily (business as usual)",
    "  - Added a new non-static connector",
    "  - Changed monitoring/alerting",
])

# Slide 17: Health Metrics
add_slide("Key Health Metrics to Watch", [
    "COVERAGE:",
    "  emails linked to bnfpvid: 103,430 (13.6% of emails)",
    "  NPIs linked to bnfpvid: 36,443 (10.9% of NPIs)",
    "  Tier1 / Tier2 ratio: 770K known vs 2.75M anonymous",
    "",
    "QUALITY:",
    "  shared_workstations flagged: 1,415 clusters (is_shared_workstation=TRUE)",
    "  is_suspicious flagged: 27 clusters (health score < 60)",
    "  bot-flagged clusters retained: 2,221 (is_bot=TRUE)",
    "  HCP clusters: 340,643 (is_hcp=TRUE)",
    "",
    "PERFORMANCE:",
    "  Full rebuild: ~4.5 hours, ~36 GB peak memory",
    "  Incremental refresh (lookback 7): ~25-45 min, ~4 GB peak",
    "",
    "STABILITY:",
    "  bn_id_persistence rows: 4,730 (active redirects)",
    "  bn_id_identity_changes: 7,958 historical merge/split events",
])

# Slide 18: What's Next
add_slide("Roadmap & Future Improvements", [
    "NEAR TERM:",
    "  - Atomic publish (view-based cutover for zero mixed-generation risk)",
    "  - Post-stitch cluster health alerting (auto-flag suspicious clusters)",
    "  - Incremental parity checks (detect drift between incremental and full rebuild)",
    "",
    "MEDIUM TERM:",
    "  - Real-time identity query API (sub-10ms lookup for downstream systems)",
    "  - Replay simulator (test rule changes before production)",
    "  - Name + geography connector (probabilistic, high risk, high reward)",
    "",
    "LONG TERM:",
    "  - Edge compaction (reduce hub storage from 9 GB to <1 GB)",
    "  - Connector parallelism (run independent connectors simultaneously)",
    "  - Cross-device ML model (learn device ownership patterns)",
])

# ============================================================================
# ANALYSIS QUERIES SECTION
# ============================================================================

add_section_slide("Analysis & Actionable Queries")

# Slide 19: Cluster Composition
add_slide("Query: Cluster Composition Overview", [
    "-- How many identifiers does each person have, by type?",
    "WITH per_bn AS (",
    "  SELECT bn_id, identifier_type, COUNT(*) AS cnt",
    "  FROM identity_hub_data.bn_id_xref",
    "  GROUP BY bn_id, identifier_type",
    ")",
    "SELECT identifier_type,",
    "  COUNT(DISTINCT bn_id) AS clusters,",
    "  ROUND(AVG(cnt), 1) AS avg_per_cluster,",
    "  MAX(cnt) AS max_per_cluster",
    "FROM per_bn GROUP BY 1 ORDER BY clusters DESC;",
    "",
    "USE CASE: Understand which identifier types are most common,",
    "which clusters are over-connected (potential false merges).",
])

# Slide 20: Find Suspicious Clusters
add_slide("Query: Suspicious Clusters Deep Dive", [
    "-- Clusters with low health scores, for investigation:",
    "WITH per_bn AS (",
    "  SELECT bn_id,",
    "    ANY_VALUE(cluster_health_score) AS health,",
    "    ANY_VALUE(is_suspicious) AS suspicious,",
    "    ANY_VALUE(is_shared_workstation) AS shared_ws,",
    "    COUNT(*) AS identifier_count",
    "  FROM identity_hub_data.bn_id_xref GROUP BY bn_id",
    ")",
    "SELECT * FROM per_bn",
    "WHERE suspicious = TRUE",
    "ORDER BY health ASC LIMIT 50;",
    "",
    "NEXT STEP: For each suspicious bn_id, check:",
    "  SELECT * FROM identity_hub_data.bn_id_xref WHERE bn_id = '...';",
])

# Slide 21: HCP Coverage
add_slide("Query: HCP Coverage Analysis", [
    "-- Which HCPs have full identity coverage?",
    "WITH hcp AS (",
    "  SELECT bn_id FROM identity_hub_data.bn_id_xref",
    "  WHERE identifier_type = 'npi_number'",
    "),",
    "hcp_detail AS (",
    "  SELECT x.bn_id,",
    "    COUNTIF(x.identifier_type = 'email') > 0 AS has_email,",
    "    COUNTIF(x.identifier_type = 'bnfpvid') > 0 AS has_browser,",
    "    COUNTIF(x.identifier_type = 'mc_euid') > 0 AS has_mailchimp",
    "  FROM identity_hub_data.bn_id_xref x",
    "  JOIN hcp ON x.bn_id = hcp.bn_id GROUP BY 1",
    ")",
    "SELECT",
    "  COUNT(*) AS total_hcps,",
    "  COUNTIF(has_email) AS with_email,",
    "  COUNTIF(has_browser) AS with_browser,",
    "  COUNTIF(has_email AND has_browser) AS full_coverage",
    "FROM hcp_detail;",
])

# Slide 22: Edge Strength Distribution
add_slide("Query: Edge Confidence Distribution", [
    "-- What does the confidence distribution look like by match rule?",
    "SELECT match_rule,",
    "  COUNT(*) AS edges,",
    "  ROUND(AVG(confidence), 3) AS avg_confidence,",
    "  COUNTIF(confidence >= 0.80) AS above_threshold,",
    "  COUNTIF(confidence < 0.80) AS below_threshold,",
    "  ROUND(SAFE_DIVIDE(COUNTIF(confidence >= 0.80),",
    "    COUNT(*)) * 100, 1) AS pct_stitchable",
    "FROM identity_hub_data.bn_id_hub",
    "GROUP BY 1 ORDER BY edges DESC;",
    "",
    "USE CASE: Monitor edge quality over time.",
    "If 'above_threshold' drops suddenly, investigate data source health.",
])

# Slide 23: Audience Segmentation Queries
add_slide("Query: Audience Segmentation", [
    "-- Known users with email who visited in last 30 days:",
    "SELECT DISTINCT x1.identifier_value AS email",
    "FROM identity_hub_data.bn_id_xref x1",
    "JOIN identity_hub_data.bn_id_xref x2 ON x1.bn_id = x2.bn_id",
    "WHERE x1.identifier_type = 'email'",
    "  AND x2.identifier_type = 'bnfpvid'",
    "  AND x2.last_seen >= DATE_SUB(CURRENT_DATE(), INTERVAL 30 DAY);",
    "",
    "-- HCPs by specialty who engaged with Mailchimp:",
    "SELECT x_npi.identifier_value AS npi, n.specialty,",
    "  x_email.identifier_value AS email",
    "FROM identity_hub_data.bn_id_xref x_npi",
    "JOIN npi_data.v_npi_provider_list n",
    "  ON x_npi.identifier_value = CAST(n.npi AS STRING)",
    "JOIN identity_hub_data.bn_id_xref x_email",
    "  ON x_npi.bn_id = x_email.bn_id AND x_email.identifier_type = 'email'",
    "WHERE x_npi.identifier_type = 'npi_number';",
])

# Slide 24: Identity Changes Monitoring
add_slide("Query: Monitor Identity Changes", [
    "-- Recent merges and splits (ops dashboard):",
    "SELECT event_date AS dt,",
    "  event_type,",
    "  COUNT(*) AS changes",
    "FROM identity_hub_data.bn_id_identity_changes",
    "GROUP BY 1, 2 ORDER BY 1 DESC LIMIT 30;",
    "",
    "-- Largest recent merges (may need investigation):",
    "SELECT *",
    "FROM identity_hub_data.bn_id_identity_changes",
    "WHERE event_type = 'MERGE'",
    "ORDER BY event_date DESC, node_count_moved DESC LIMIT 20;",
    "",
    "-- bn_id stability: how many redirects exist?",
    "SELECT COUNT(*) AS redirects,",
    "  COUNT(DISTINCT old_bn_id) AS old_ids,",
    "  COUNT(DISTINCT current_bn_id) AS current_ids",
    "FROM identity_hub_data.bn_id_persistence;",
])

# Slide 25: Compound Cross-Device Analysis
add_slide("Query: Cross-Device Signal Analysis", [
    "-- How effective is the device fingerprint connector?",
    "SELECT",
    "  match_rule,",
    "  COUNT(*) AS total_edges,",
    "  COUNTIF(effective_confidence >= 0.80) AS stitching_edges,",
    "  COUNTIF(effective_confidence < 0.80) AS audit_only_edges,",
    "  ROUND(AVG(confidence), 3) AS avg_confidence",
    "FROM identity_hub_data.bn_id_hub",
    "WHERE match_rule IN ('DEVICE_STAT_ID', 'COMPOUND_DEVICE_IP_TIME')",
    "GROUP BY match_rule;",
    "",
    "-- Which clusters are linked via device fingerprint?",
    "SELECT h.bn_id, COUNT(*) AS device_edges",
    "FROM identity_hub_data.bn_id_hub h",
    "WHERE h.match_rule IN ('DEVICE_STAT_ID', 'COMPOUND_DEVICE_IP_TIME')",
    "  AND h.effective_confidence >= 0.80",
    "GROUP BY h.bn_id ORDER BY device_edges DESC LIMIT 20;",
])

# Slide 26: Shared Workstation Analysis
add_slide("Query: Shared Workstation Analysis", [
    "-- Clusters flagged as shared workstations:",
    "WITH per_bn AS (",
    "  SELECT bn_id, ANY_VALUE(is_shared_workstation) AS ws,",
    "    COUNT(*) AS identifiers",
    "  FROM identity_hub_data.bn_id_xref GROUP BY bn_id",
    ")",
    "SELECT COUNTIF(ws) AS ws_clusters,",
    "  ROUND(AVG(IF(ws, identifiers, NULL)), 1) AS avg_ids_in_ws,",
    "  MAX(IF(ws, identifiers, NULL)) AS max_ids_in_ws",
    "FROM per_bn;",
    "",
    "-- Which shared workstations have the most users?",
    "SELECT bn_id, COUNT(*) AS identifiers,",
    "  COUNTIF(identifier_type = 'email') AS emails,",
    "  COUNTIF(identifier_type = 'mc_euid') AS mc_euids",
    "FROM identity_hub_data.bn_id_xref",
    "WHERE is_shared_workstation = TRUE",
    "GROUP BY 1 ORDER BY emails DESC LIMIT 20;",
])

# Save
output_path = 'docs/Identity_Hub_Presentation_v3.pptx'
prs.save(output_path)
print(f"Saved to {output_path}")
