"""Generate PowerPoint presentation for LimeSurvey Pipeline"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# Create presentation with 16:9 aspect ratio
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Color scheme
BLUE = RGBColor(0x1a, 0x73, 0xe8)
GREEN = RGBColor(0x34, 0xa8, 0x53)
RED = RGBColor(0xea, 0x43, 0x35)
ORANGE = RGBColor(0xfb, 0xbc, 0x04)
PURPLE = RGBColor(0x9c, 0x27, 0xb0)
GRAY = RGBColor(0x60, 0x7d, 0x8b)
DARK = RGBColor(0x33, 0x33, 0x33)
LIGHT_BLUE = RGBColor(0xe8, 0xf0, 0xfe)

def add_title_slide(prs, title, subtitle=""):
    """Add a title slide"""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(12.333), Inches(1.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(54)
    p.font.bold = True
    p.font.color.rgb = BLUE
    p.alignment = PP_ALIGN.CENTER

    if subtitle:
        sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(4), Inches(12.333), Inches(0.75))
        tf = sub_box.text_frame
        p = tf.paragraphs[0]
        p.text = subtitle
        p.font.size = Pt(28)
        p.font.color.rgb = RGBColor(0x66, 0x66, 0x66)
        p.alignment = PP_ALIGN.CENTER

    return slide

def add_content_slide(prs, title, content_func):
    """Add a content slide with title"""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)

    # Title bar
    title_shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(1.2))
    title_shape.fill.solid()
    title_shape.fill.fore_color.rgb = BLUE
    title_shape.line.fill.background()

    # Title text
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0xff, 0xff, 0xff)

    # Add content
    content_func(slide)

    return slide

def add_text_box(slide, left, top, width, height, text, font_size=18, bold=False, color=DARK, align=PP_ALIGN.LEFT):
    """Helper to add a text box"""
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.alignment = align
    return box

def add_bullet_list(slide, left, top, width, height, items, font_size=18):
    """Helper to add a bullet list"""
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True

    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"• {item}"
        p.font.size = Pt(font_size)
        p.font.color.rgb = DARK
        p.space_after = Pt(8)
    return box

def add_stat_box(slide, left, top, number, label):
    """Add a statistic box"""
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(2), Inches(1.2))
    shape.fill.solid()
    shape.fill.fore_color.rgb = BLUE
    shape.line.fill.background()

    # Number
    num_box = slide.shapes.add_textbox(Inches(left), Inches(top + 0.15), Inches(2), Inches(0.6))
    tf = num_box.text_frame
    p = tf.paragraphs[0]
    p.text = number
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0xff, 0xff, 0xff)
    p.alignment = PP_ALIGN.CENTER

    # Label
    label_box = slide.shapes.add_textbox(Inches(left), Inches(top + 0.7), Inches(2), Inches(0.4))
    tf = label_box.text_frame
    p = tf.paragraphs[0]
    p.text = label
    p.font.size = Pt(14)
    p.font.color.rgb = RGBColor(0xff, 0xff, 0xff)
    p.alignment = PP_ALIGN.CENTER

def add_layer_box(slide, left, top, num, name, confidence, desc, color):
    """Add a matching layer box"""
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(12), Inches(0.8))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()

    # Number
    num_box = slide.shapes.add_textbox(Inches(left + 0.2), Inches(top + 0.15), Inches(0.5), Inches(0.5))
    tf = num_box.text_frame
    p = tf.paragraphs[0]
    p.text = str(num)
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0xff, 0xff, 0xff)

    # Name
    name_box = slide.shapes.add_textbox(Inches(left + 0.8), Inches(top + 0.1), Inches(4), Inches(0.35))
    tf = name_box.text_frame
    p = tf.paragraphs[0]
    p.text = name
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0xff, 0xff, 0xff)

    # Confidence
    conf_box = slide.shapes.add_textbox(Inches(left + 9), Inches(top + 0.1), Inches(2.5), Inches(0.35))
    tf = conf_box.text_frame
    p = tf.paragraphs[0]
    p.text = confidence
    p.font.size = Pt(14)
    p.font.color.rgb = RGBColor(0xff, 0xff, 0xff)
    p.alignment = PP_ALIGN.RIGHT

    # Description
    desc_box = slide.shapes.add_textbox(Inches(left + 0.8), Inches(top + 0.45), Inches(10), Inches(0.35))
    tf = desc_box.text_frame
    p = tf.paragraphs[0]
    p.text = desc
    p.font.size = Pt(12)
    p.font.color.rgb = RGBColor(0xff, 0xff, 0xff)

# SLIDE 1: Title
slide = add_title_slide(prs, "LimeSurvey Data Pipeline", "Transforming Survey Chaos into Actionable Insights")
add_stat_box(slide, 3.5, 5, "163", "Surveys")
add_stat_box(slide, 5.7, 5, "6,000+", "Questions")
add_stat_box(slide, 7.9, 5, "1M+", "Responses")
add_text_box(slide, 0.5, 6.5, 12.333, 0.5, "BioNews Data Engineering", 16, color=RGBColor(0x99, 0x99, 0x99), align=PP_ALIGN.CENTER)

# SLIDE 2: The Problem
def slide2_content(slide):
    add_text_box(slide, 0.5, 1.5, 12, 0.5, "The same question is asked in many different ways across 163 surveys:", 20)

    # Table header
    add_text_box(slide, 0.5, 2.2, 4, 0.4, "Survey A", 16, bold=True, color=RGBColor(0xff, 0xff, 0xff))
    add_text_box(slide, 4.5, 2.2, 4, 0.4, "Survey B", 16, bold=True, color=RGBColor(0xff, 0xff, 0xff))
    add_text_box(slide, 8.5, 2.2, 4, 0.4, "Survey C", 16, bold=True, color=RGBColor(0xff, 0xff, 0xff))

    for i, left in enumerate([0.5, 4.5, 8.5]):
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(left), Inches(2.1), Inches(3.8), Inches(0.5))
        shape.fill.solid()
        shape.fill.fore_color.rgb = BLUE
        shape.line.fill.background()

    # Row 1
    add_text_box(slide, 0.5, 2.7, 4, 0.4, '"What is your age?"', 16)
    add_text_box(slide, 4.5, 2.7, 4, 0.4, '"How old are you?"', 16)
    add_text_box(slide, 8.5, 2.7, 4, 0.4, '"Please select your age range"', 16)

    # Row 2
    add_text_box(slide, 0.5, 3.2, 4, 0.4, '"18-24"', 16)
    add_text_box(slide, 4.5, 3.2, 4, 0.4, '"18 to 24 years"', 16)
    add_text_box(slide, 8.5, 3.2, 4, 0.4, '"Young adult (18-24)"', 16)

    # Warning box
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(4.2), Inches(12), Inches(1.2))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0xff, 0xf3, 0xe0)
    shape.line.color.rgb = ORANGE

    add_text_box(slide, 0.7, 4.4, 11.5, 1, 'The Challenge: A BI analyst cannot easily answer "How many patients are aged 18-24?" because the question and answers are phrased differently in each survey.', 18, bold=False, color=DARK)

add_content_slide(prs, "The Problem: Survey Data Chaos", slide2_content)

# SLIDE 3: Solution Overview
def slide3_content(slide):
    stages = [
        ("Stage 1: EXTRACTION", "Pull raw data from LimeSurvey MySQL database", LIGHT_BLUE),
        ("Stage 2: ENRICHMENT", "Convert codes to readable text", LIGHT_BLUE),
        ("Stage 3: STANDARDIZATION", "AI-powered question & answer matching", RGBColor(0xe8, 0xf5, 0xe9)),
        ("Stage 4: ANALYTICS OUTPUT", "Create BI-ready wide tables", LIGHT_BLUE),
    ]

    top = 1.6
    for title, desc, color in stages:
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(2), Inches(top), Inches(9), Inches(0.9))
        shape.fill.solid()
        shape.fill.fore_color.rgb = color
        shape.line.color.rgb = BLUE if color == LIGHT_BLUE else GREEN

        add_text_box(slide, 2.2, top + 0.15, 8.5, 0.35, title, 20, bold=True, color=DARK, align=PP_ALIGN.CENTER)
        add_text_box(slide, 2.2, top + 0.5, 8.5, 0.35, desc, 14, color=RGBColor(0x66, 0x66, 0x66), align=PP_ALIGN.CENTER)

        top += 1.0
        if top < 5.5:
            add_text_box(slide, 6.3, top - 0.1, 0.5, 0.3, "↓", 28, color=BLUE, align=PP_ALIGN.CENTER)
            top += 0.3

add_content_slide(prs, "The Solution: 4-Stage Pipeline", slide3_content)

# SLIDE 4: Stage 1 - Extraction
def slide4_content(slide):
    add_text_box(slide, 0.5, 1.5, 6, 0.5, "What Happens", 24, bold=True, color=BLUE)
    add_bullet_list(slide, 0.5, 2.1, 6, 2.5, [
        "Pipeline connects to LimeSurvey via secure SSH tunnel",
        "Copies survey responses from MySQL to BigQuery",
        "Runs automatically every night",
        'Data arrives in "long format" (one row per answer)'
    ], 18)

    add_text_box(slide, 7, 1.5, 5.5, 0.5, "Technologies", 24, bold=True, color=BLUE)

    # Tech badges
    techs = ["Python", "SSH Tunnel", "MySQL", "BigQuery"]
    left = 7
    for tech in techs:
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(2.1), Inches(1.3), Inches(0.4))
        shape.fill.solid()
        shape.fill.fore_color.rgb = GREEN
        shape.line.fill.background()
        add_text_box(slide, left + 0.1, 2.15, 1.1, 0.3, tech, 12, color=RGBColor(0xff, 0xff, 0xff), align=PP_ALIGN.CENTER)
        left += 1.4

    # Output box
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7), Inches(3), Inches(5.5), Inches(1.5))
    shape.fill.solid()
    shape.fill.fore_color.rgb = LIGHT_BLUE
    shape.line.color.rgb = BLUE

    add_text_box(slide, 7.2, 3.2, 5, 0.4, "Output Table:", 18, bold=True, color=DARK)
    add_text_box(slide, 7.2, 3.6, 5, 0.4, "lime_surveys_columnar", 16, color=BLUE)
    add_text_box(slide, 7.2, 4, 5, 0.4, "~1 million rows of raw data", 14, color=RGBColor(0x66, 0x66, 0x66))

add_content_slide(prs, "Stage 1: Extraction", slide4_content)

# SLIDE 5: Stage 2 - Enrichment
def slide5_content(slide):
    # Before box
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(1.5), Inches(5.5), Inches(2))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0xfc, 0xe4, 0xec)
    shape.line.color.rgb = RGBColor(0xe9, 0x1e, 0x63)

    add_text_box(slide, 0.7, 1.7, 5, 0.4, "Before (Raw)", 22, bold=True, color=RGBColor(0xc2, 0x18, 0x5b))
    add_text_box(slide, 0.7, 2.2, 5, 0.3, "Question ID: Q001", 16)
    add_text_box(slide, 0.7, 2.6, 5, 0.3, "Answer Code: AO01", 16)
    add_text_box(slide, 0.7, 3, 5, 0.3, "Not human-readable!", 16, color=RGBColor(0xc2, 0x18, 0x5b))

    # Arrow
    add_text_box(slide, 6.2, 2.2, 0.8, 0.5, "→", 48, color=BLUE, align=PP_ALIGN.CENTER)

    # After box
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7), Inches(1.5), Inches(5.5), Inches(2))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0xe8, 0xf5, 0xe9)
    shape.line.color.rgb = RGBColor(0x4c, 0xaf, 0x50)

    add_text_box(slide, 7.2, 1.7, 5, 0.4, "After (Enriched)", 22, bold=True, color=RGBColor(0x38, 0x8e, 0x3c))
    add_text_box(slide, 7.2, 2.2, 5, 0.3, 'Question: "What is your age?"', 16)
    add_text_box(slide, 7.2, 2.6, 5, 0.3, 'Answer: "18-24 years old"', 16)
    add_text_box(slide, 7.2, 3, 5, 0.3, "Now we understand it!", 16, color=RGBColor(0x38, 0x8e, 0x3c))

    # Key actions
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(4), Inches(12), Inches(2.5))
    shape.fill.solid()
    shape.fill.fore_color.rgb = LIGHT_BLUE
    shape.line.color.rgb = BLUE

    add_text_box(slide, 0.7, 4.2, 11.5, 0.4, "Key Actions", 20, bold=True, color=BLUE)
    add_bullet_list(slide, 0.7, 4.7, 11.5, 1.5, [
        "Looks up actual question text for each question ID",
        "Converts answer codes to human-readable text",
        "Removes HTML formatting",
        "Adds rows to permanent production table"
    ], 16)

add_content_slide(prs, "Stage 2: Enrichment", slide5_content)

# SLIDE 6: Stage 3 Overview
def slide6_content(slide):
    add_text_box(slide, 0.5, 1.5, 12, 0.6, 'The "smart matching" process uses 6 layers of artificial intelligence to find the canonical name for each question.', 20)

    # Input column
    add_text_box(slide, 0.5, 2.3, 5.5, 0.5, "Input (Variations)", 22, bold=True, color=BLUE)
    add_bullet_list(slide, 0.5, 2.9, 5.5, 2, [
        '"What is your age?"',
        '"How old are you?"',
        '"Please select your age range"',
        '"Whats your age"'
    ], 18)

    # Arrow
    add_text_box(slide, 6.2, 3.5, 0.8, 0.5, "→", 48, color=BLUE, align=PP_ALIGN.CENTER)

    # Output column
    add_text_box(slide, 7, 2.3, 5.5, 0.5, "Output (Standardized)", 22, bold=True, color=BLUE)

    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7), Inches(2.9), Inches(5.5), Inches(1.8))
    shape.fill.solid()
    shape.fill.fore_color.rgb = LIGHT_BLUE
    shape.line.color.rgb = BLUE

    add_text_box(slide, 7.2, 3.2, 5, 0.8, "age", 48, bold=True, color=BLUE, align=PP_ALIGN.CENTER)
    add_text_box(slide, 7.2, 4, 5, 0.4, "One canonical name for all variations", 14, color=RGBColor(0x66, 0x66, 0x66), align=PP_ALIGN.CENTER)

    # Tech badges
    techs = ["Sentence Transformers", "UMLS Medical Terms", "spaCy NER", "Machine Learning", "RapidFuzz"]
    left = 0.5
    for tech in techs:
        w = len(tech) * 0.12 + 0.3
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(5.5), Inches(w), Inches(0.4))
        shape.fill.solid()
        shape.fill.fore_color.rgb = GREEN
        shape.line.fill.background()
        add_text_box(slide, left + 0.1, 5.55, w - 0.2, 0.3, tech, 12, color=RGBColor(0xff, 0xff, 0xff), align=PP_ALIGN.CENTER)
        left += w + 0.15

add_content_slide(prs, "Stage 3: AI-Powered Standardization", slide6_content)

# SLIDE 7: The 6 Matching Layers
def slide7_content(slide):
    layers = [
        (1, "EXACT MATCH", "100% confidence", "If we've seen this exact question before, use the same mapping", BLUE),
        (2, "MACHINE LEARNING", "90-99% confidence", "Trained model predicts category based on patterns from thousands of mappings", GREEN),
        (3, "SEMANTIC SIMILARITY", "85-99% confidence", '"How old are you?" matches "What is your age?" - same meaning!', ORANGE),
        (4, "MEDICAL TERMS (UMLS)", "85-98% confidence", '"ALS" matches "Amyotrophic Lateral Sclerosis" via medical database', RED),
        (5, "ENTITY RECOGNITION", "80-93% confidence", "Identifies diseases, chemicals, body parts in questions using spaCy", PURPLE),
        (6, "FUZZY MATCHING", "70-85% confidence", '"Whats your age" still matches "What is your age?" - catches typos!', GRAY),
    ]

    top = 1.5
    for num, name, conf, desc, color in layers:
        add_layer_box(slide, 0.5, top, num, name, conf, desc, color)
        top += 0.9

add_content_slide(prs, "The 6 Matching Layers", slide7_content)

# SLIDE 8: How Matching Works
def slide8_content(slide):
    diagram = '''New Question: "How old are you?"
           |
           v
    +------------------+
    | Layer 1: Exact   |---> No match found
    +------------------+
           |
           v
    +------------------+
    | Layer 2: ML      |---> 92% confidence for "age"
    +------------------+
           |
           v
    +------------------+
    | Layer 3: Semantic|---> 97% confidence for "age"  <-- STOP! Above 95%
    +------------------+

    Result: "How old are you?" --> standardized_question = "age"'''

    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(1.5), Inches(12), Inches(4))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0xf8, 0xf9, 0xfa)
    shape.line.color.rgb = RGBColor(0xdd, 0xdd, 0xdd)

    box = slide.shapes.add_textbox(Inches(0.7), Inches(1.7), Inches(11.5), Inches(3.5))
    tf = box.text_frame
    p = tf.paragraphs[0]
    p.text = diagram
    p.font.size = Pt(14)
    p.font.name = "Courier New"
    p.font.color.rgb = DARK

    # Key insight box
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(5.7), Inches(12), Inches(1))
    shape.fill.solid()
    shape.fill.fore_color.rgb = LIGHT_BLUE
    shape.line.color.rgb = BLUE

    add_text_box(slide, 0.7, 5.9, 11.5, 0.7, "Key Insight: The system stops at the first layer that gives 95%+ confidence, saving processing time and ensuring high accuracy.", 18, color=DARK)

add_content_slide(prs, "How Matching Works", slide8_content)

# SLIDE 9: Stage 4 - Analytics Output
def slide9_content(slide):
    add_text_box(slide, 0.5, 1.5, 12, 0.5, "Transform from Long to Wide Format", 22, bold=True, color=BLUE)

    # Before table
    add_text_box(slide, 0.5, 2.1, 5.5, 0.4, "BEFORE (Long Format)", 18, bold=True)

    # Simple table representation
    headers = ["response_id", "question", "answer"]
    rows = [["12345", "age", "25-34"], ["12345", "gender", "Female"], ["12345", "diagnosis", "ALS"]]

    top = 2.6
    for i, h in enumerate(headers):
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5 + i*1.8), Inches(top), Inches(1.7), Inches(0.35))
        shape.fill.solid()
        shape.fill.fore_color.rgb = BLUE
        shape.line.fill.background()
        add_text_box(slide, 0.55 + i*1.8, top + 0.05, 1.6, 0.25, h, 11, bold=True, color=RGBColor(0xff, 0xff, 0xff), align=PP_ALIGN.CENTER)

    top += 0.4
    for row in rows:
        for i, cell in enumerate(row):
            add_text_box(slide, 0.55 + i*1.8, top, 1.6, 0.3, cell, 11, align=PP_ALIGN.CENTER)
        top += 0.35

    add_text_box(slide, 0.5, 4.2, 5.5, 0.3, "One row per answer", 12, color=RGBColor(0x66, 0x66, 0x66))

    # Arrow
    add_text_box(slide, 6, 3, 0.8, 0.5, "→", 36, color=BLUE, align=PP_ALIGN.CENTER)

    # After table
    add_text_box(slide, 6.8, 2.1, 6, 0.4, "AFTER (Wide Format)", 18, bold=True)

    headers2 = ["response_id", "age", "gender", "diagnosis"]
    rows2 = [["12345", "25-34", "Female", "ALS"], ["12346", "45-54", "Male", "MS"]]

    top = 2.6
    for i, h in enumerate(headers2):
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(6.8 + i*1.5), Inches(top), Inches(1.4), Inches(0.35))
        shape.fill.solid()
        shape.fill.fore_color.rgb = BLUE
        shape.line.fill.background()
        add_text_box(slide, 6.85 + i*1.5, top + 0.05, 1.3, 0.25, h, 11, bold=True, color=RGBColor(0xff, 0xff, 0xff), align=PP_ALIGN.CENTER)

    top += 0.4
    for row in rows2:
        for i, cell in enumerate(row):
            add_text_box(slide, 6.85 + i*1.5, top, 1.3, 0.3, cell, 11, align=PP_ALIGN.CENTER)
        top += 0.35

    add_text_box(slide, 6.8, 3.8, 6, 0.3, "One row per respondent, ~148 columns", 12, color=RGBColor(0x66, 0x66, 0x66))

    # Output box
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(4.8), Inches(12), Inches(1.2))
    shape.fill.solid()
    shape.fill.fore_color.rgb = LIGHT_BLUE
    shape.line.color.rgb = BLUE

    add_text_box(slide, 0.7, 5, 11.5, 0.8, "Output: lime_surveys_wide - Ready for Looker dashboards and SQL queries!", 20, bold=True, color=DARK)

add_content_slide(prs, "Stage 4: Analytics Output", slide9_content)

# SLIDE 10: Lookup Tables
def slide10_content(slide):
    add_text_box(slide, 0.5, 1.5, 12, 0.5, 'For standardization to work, we need reference lists of "correct" values:', 18)

    tables = [
        ("standardized_questions_lookup", "148 canonical question names", '"age", "gender", "diagnosis_confirmed"'),
        ("standardized_answers_lookup", "Approved answers per question", 'For "gender": male, female, non_binary'),
        ("categories_lookup", "Question categories", "Demographics, Medical History, Treatment"),
        ("question_mapping", "Raw to standardized questions", '"What is your age?" → "age"'),
        ("answer_mapping", "Raw to standardized answers", '"18-24 years" → "18_to_24"'),
    ]

    # Header
    top = 2.2
    headers = ["Table", "Purpose", "Example"]
    widths = [4, 3.5, 4]
    left = 0.5
    for i, h in enumerate(headers):
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(left), Inches(top), Inches(widths[i]), Inches(0.4))
        shape.fill.solid()
        shape.fill.fore_color.rgb = BLUE
        shape.line.fill.background()
        add_text_box(slide, left + 0.1, top + 0.05, widths[i] - 0.2, 0.3, h, 14, bold=True, color=RGBColor(0xff, 0xff, 0xff))
        left += widths[i] + 0.1

    top += 0.5
    for table, purpose, example in tables:
        left = 0.5
        add_text_box(slide, left + 0.1, top, 3.8, 0.4, table, 13, color=BLUE)
        add_text_box(slide, left + 4.1, top, 3.3, 0.4, purpose, 13)
        add_text_box(slide, left + 7.7, top, 3.8, 0.4, example, 13)
        top += 0.5

add_content_slide(prs, "Supporting Lookup Tables", slide10_content)

# SLIDE 11: Manual Review
def slide11_content(slide):
    add_text_box(slide, 0.5, 1.5, 12, 0.5, "Not every question can be matched automatically. When the AI is unsure:", 18)

    steps = [
        ("1. Flag", "Unmapped question appears in Streamlit UI", LIGHT_BLUE),
        ("2. Review", "Human sees raw question + suggested matches", LIGHT_BLUE),
        ("3. Select", "Choose correct standardized question", LIGHT_BLUE),
        ("4. Learn", "Mapping saved for future matches", RGBColor(0xe8, 0xf5, 0xe9)),
    ]

    left = 0.5
    for title, desc, color in steps:
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(2.3), Inches(2.5), Inches(1.2))
        shape.fill.solid()
        shape.fill.fore_color.rgb = color
        shape.line.color.rgb = BLUE if color == LIGHT_BLUE else GREEN

        add_text_box(slide, left + 0.1, 2.5, 2.3, 0.4, title, 18, bold=True, color=DARK, align=PP_ALIGN.CENTER)
        add_text_box(slide, left + 0.1, 2.9, 2.3, 0.5, desc, 11, color=RGBColor(0x66, 0x66, 0x66), align=PP_ALIGN.CENTER)

        if left < 9:
            add_text_box(slide, left + 2.7, 2.7, 0.5, 0.4, "→", 24, color=BLUE, align=PP_ALIGN.CENTER)
        left += 3.1

    # Streamlit tools box
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(4), Inches(12), Inches(2))
    shape.fill.solid()
    shape.fill.fore_color.rgb = LIGHT_BLUE
    shape.line.color.rgb = BLUE

    add_text_box(slide, 0.7, 4.2, 11.5, 0.4, "Streamlit UI Tools", 20, bold=True, color=BLUE)
    add_bullet_list(slide, 0.7, 4.7, 11.5, 1, [
        "limesurvey_mapper.py - Map questions to standardized names",
        "answer_standardizer.py - Map answers to standardized values"
    ], 16)

add_content_slide(prs, "Manual Review Process", slide11_content)

# SLIDE 12: Key Tables Summary
def slide12_content(slide):
    tables = [
        ("lime_surveys_columnar", "Staging", "~1M", "Can rebuild"),
        ("lime_surveys_columnar_completed", "SOURCE OF TRUTH", "~1.3M", "NEVER DELETE"),
        ("lime_surveys_wide", "Analytics", "~65K", "Safe to rebuild"),
        ("respondent_profiles", "BI Output", "~2.8M", "Safe to rebuild"),
        ("ml_question_model", "ML Storage", "1", "No"),
    ]

    # Header
    top = 1.6
    headers = ["Table", "Type", "Rows", "Can Delete?"]
    widths = [5, 2.5, 1.5, 2.5]
    left = 0.5
    for i, h in enumerate(headers):
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(left), Inches(top), Inches(widths[i]), Inches(0.4))
        shape.fill.solid()
        shape.fill.fore_color.rgb = BLUE
        shape.line.fill.background()
        add_text_box(slide, left + 0.1, top + 0.05, widths[i] - 0.2, 0.3, h, 14, bold=True, color=RGBColor(0xff, 0xff, 0xff))
        left += widths[i] + 0.1

    top += 0.5
    for table, ttype, rows, delete in tables:
        left = 0.5
        is_critical = "NEVER" in delete

        if is_critical:
            shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(top - 0.05), Inches(11.6), Inches(0.45))
            shape.fill.solid()
            shape.fill.fore_color.rgb = RGBColor(0xff, 0xf3, 0xe0)
            shape.line.fill.background()

        add_text_box(slide, left + 0.1, top, 4.8, 0.4, table, 14, color=BLUE)
        add_text_box(slide, left + 5.1, top, 2.3, 0.4, ttype, 14, bold=is_critical)
        add_text_box(slide, left + 7.7, top, 1.3, 0.4, rows, 14)
        add_text_box(slide, left + 9.3, top, 2.3, 0.4, delete, 14, bold=is_critical, color=RED if is_critical else DARK)
        top += 0.5

    # Warning
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(4.5), Inches(12), Inches(1))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0xff, 0xf3, 0xe0)
    shape.line.color.rgb = ORANGE

    add_text_box(slide, 0.7, 4.7, 11.5, 0.7, "Critical: lime_surveys_columnar_completed is APPEND-ONLY. All standardization mappings are stored here. Never rebuild or drop!", 16, color=DARK)

add_content_slide(prs, "Key Tables Summary", slide12_content)

# SLIDE 13: Data Flow
def slide13_content(slide):
    diagram = '''MySQL (LimeSurvey)
    |
    +-- lime_surveys, lime_questions, lime_answers, etc.
    |   (extracted via SSH tunnel every night)
    |
    +-- lime_survey_XXXXX (individual response tables)
            |
            v
    lime_surveys_columnar (STAGING)
            |
            | sql/limesurvey_enrich_columnar.sql
            v
    lime_surveys_columnar_completed (SOURCE OF TRUTH)
            |
            +-- limesurvey_process_columnar.py (NLP standardization)
            |
            +-- limesurvey_build_wide_table.py
            |       |
            |       v
            |   lime_surveys_wide (ANALYTICS) --> Looker Dashboards
            |
            +-- survey_profile_builder.py
                    |
                    v
                respondent_profiles (BI)'''

    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(1.5), Inches(12), Inches(5.2))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0xf8, 0xf9, 0xfa)
    shape.line.color.rgb = RGBColor(0xdd, 0xdd, 0xdd)

    box = slide.shapes.add_textbox(Inches(0.7), Inches(1.7), Inches(11.5), Inches(5))
    tf = box.text_frame
    p = tf.paragraphs[0]
    p.text = diagram
    p.font.size = Pt(13)
    p.font.name = "Courier New"
    p.font.color.rgb = DARK

add_content_slide(prs, "Complete Data Flow", slide13_content)

# SLIDE 14: Why This Matters
def slide14_content(slide):
    # Before box
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(1.5), Inches(5.5), Inches(3.5))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0xfc, 0xe4, 0xec)
    shape.line.color.rgb = RGBColor(0xe9, 0x1e, 0x63)

    add_text_box(slide, 0.7, 1.7, 5, 0.5, "Before Standardization", 22, bold=True, color=RGBColor(0xc2, 0x18, 0x5b))
    add_text_box(slide, 0.7, 2.3, 5, 0.5, '"How many patients said they are 18-24?"', 16)
    add_text_box(slide, 0.7, 3, 5, 1, "Requires querying 50+ different surveys with 50+ different question phrasings", 16)
    add_text_box(slide, 0.7, 4.2, 5, 0.5, "Hours of manual work", 20, bold=True, color=RGBColor(0xc2, 0x18, 0x5b))

    # After box
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.5), Inches(1.5), Inches(6), Inches(3.5))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0xe8, 0xf5, 0xe9)
    shape.line.color.rgb = RGBColor(0x4c, 0xaf, 0x50)

    add_text_box(slide, 6.7, 1.7, 5.5, 0.5, "After Standardization", 22, bold=True, color=RGBColor(0x38, 0x8e, 0x3c))
    add_text_box(slide, 6.7, 2.3, 5.5, 0.4, "One simple query:", 16)

    # Code box
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.7), Inches(2.8), Inches(5.5), Inches(1.2))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0x33, 0x33, 0x33)
    shape.line.fill.background()

    code = "SELECT COUNT(*)\nFROM lime_surveys_wide\nWHERE age = '18_to_24'"
    box = slide.shapes.add_textbox(Inches(6.9), Inches(2.9), Inches(5.1), Inches(1))
    tf = box.text_frame
    p = tf.paragraphs[0]
    p.text = code
    p.font.size = Pt(14)
    p.font.name = "Courier New"
    p.font.color.rgb = RGBColor(0x4c, 0xaf, 0x50)

    add_text_box(slide, 6.7, 4.2, 5.5, 0.5, "Seconds!", 20, bold=True, color=RGBColor(0x38, 0x8e, 0x3c))

add_content_slide(prs, "Why This Matters", slide14_content)

# SLIDE 15: Technologies
def slide15_content(slide):
    # Left column
    add_text_box(slide, 0.5, 1.5, 5.5, 0.4, "Data Infrastructure", 20, bold=True, color=BLUE)
    techs1 = ["Google BigQuery", "Google Cloud Storage", "MySQL", "SSH Tunneling"]
    left = 0.5
    top = 2
    for tech in techs1:
        w = len(tech) * 0.11 + 0.3
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(w), Inches(0.35))
        shape.fill.solid()
        shape.fill.fore_color.rgb = GREEN
        shape.line.fill.background()
        add_text_box(slide, left + 0.1, top + 0.05, w - 0.2, 0.25, tech, 11, color=RGBColor(0xff, 0xff, 0xff), align=PP_ALIGN.CENTER)
        left += w + 0.1
        if left > 5.5:
            left = 0.5
            top += 0.45

    add_text_box(slide, 0.5, 3, 5.5, 0.4, "AI/ML Components", 20, bold=True, color=BLUE)
    techs2 = ["Sentence Transformers", "RandomForest", "UMLS Medical DB", "spaCy NER", "RapidFuzz"]
    left = 0.5
    top = 3.5
    for tech in techs2:
        w = len(tech) * 0.11 + 0.3
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(w), Inches(0.35))
        shape.fill.solid()
        shape.fill.fore_color.rgb = GREEN
        shape.line.fill.background()
        add_text_box(slide, left + 0.1, top + 0.05, w - 0.2, 0.25, tech, 11, color=RGBColor(0xff, 0xff, 0xff), align=PP_ALIGN.CENTER)
        left += w + 0.1
        if left > 5.5:
            left = 0.5
            top += 0.45

    # Right column
    add_text_box(slide, 7, 1.5, 5.5, 0.4, "Development Stack", 20, bold=True, color=BLUE)
    techs3 = ["Python 3.13", "Streamlit", "scikit-learn", "pandas"]
    left = 7
    top = 2
    for tech in techs3:
        w = len(tech) * 0.11 + 0.3
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(w), Inches(0.35))
        shape.fill.solid()
        shape.fill.fore_color.rgb = GREEN
        shape.line.fill.background()
        add_text_box(slide, left + 0.1, top + 0.05, w - 0.2, 0.25, tech, 11, color=RGBColor(0xff, 0xff, 0xff), align=PP_ALIGN.CENTER)
        left += w + 0.1
        if left > 12:
            left = 7
            top += 0.45

    add_text_box(slide, 7, 3, 5.5, 0.4, "ML Models", 20, bold=True, color=BLUE)
    techs4 = ["all-MiniLM-L6-v2", "en_core_sci_md", "en_ner_bc5cdr_md"]
    left = 7
    top = 3.5
    for tech in techs4:
        w = len(tech) * 0.1 + 0.3
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(w), Inches(0.35))
        shape.fill.solid()
        shape.fill.fore_color.rgb = GREEN
        shape.line.fill.background()
        add_text_box(slide, left + 0.1, top + 0.05, w - 0.2, 0.25, tech, 11, color=RGBColor(0xff, 0xff, 0xff), align=PP_ALIGN.CENTER)
        left += w + 0.1

add_content_slide(prs, "Technologies Used", slide15_content)

# SLIDE 16: Summary
slide = add_title_slide(prs, "Summary", "")

# Summary box
shape = prs.slides[-1].shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(2), Inches(2), Inches(9), Inches(3))
shape.fill.solid()
shape.fill.fore_color.rgb = LIGHT_BLUE
shape.line.color.rgb = BLUE

add_bullet_list(prs.slides[-1], 2.3, 2.3, 8.5, 2, [
    "Extract raw survey data nightly from LimeSurvey",
    "Enrich with human-readable question and answer text",
    "Standardize using 6 layers of AI matching",
    "Output BI-ready tables for analytics"
], 20)

add_text_box(prs.slides[-1], 2, 5.3, 9, 0.5, "Result: 163 surveys with 6,000+ questions unified into", 18, align=PP_ALIGN.CENTER)
add_text_box(prs.slides[-1], 2, 5.8, 9, 0.6, "148 standardized question columns", 32, bold=True, color=BLUE, align=PP_ALIGN.CENTER)

# Save
output_path = "c:/orchestrator/docs/LimeSurvey_Pipeline_Presentation.pptx"
prs.save(output_path)
print(f"Presentation saved to: {output_path}")
