"""Generate updated PowerPoint presentation for the LimeSurvey Pipeline (v2).

Covers the current state of the system:
  - Unified Streamlit mapper (other apps deprecated)
  - Blended suggestion engine (fuzzy + semantic + co-occurrence + fingerprints)
  - Manual-assignment protection, undo, audit log
  - Daily 7-day lookback + nightly retrain
  - Identity hub bridge via survey-response email (not lime_participants)

Run:
    python docs/generate_pptx_v2.py
Output:
    docs/LimeSurvey_Pipeline_v2.pptx
"""

from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# =====================================================================
# Color scheme
# =====================================================================
BLUE = RGBColor(0x1a, 0x73, 0xe8)
DARK_BLUE = RGBColor(0x0d, 0x47, 0xa1)
GREEN = RGBColor(0x34, 0xa8, 0x53)
RED = RGBColor(0xea, 0x43, 0x35)
ORANGE = RGBColor(0xfb, 0xbc, 0x04)
PURPLE = RGBColor(0x9c, 0x27, 0xb0)
GRAY = RGBColor(0x60, 0x7d, 0x8b)
DARK = RGBColor(0x33, 0x33, 0x33)
MUTED = RGBColor(0x66, 0x66, 0x66)
LIGHT_GRAY = RGBColor(0xf5, 0xf7, 0xfa)
LIGHT_BLUE = RGBColor(0xe8, 0xf0, 0xfe)
LIGHT_GREEN = RGBColor(0xe8, 0xf5, 0xe9)
LIGHT_AMBER = RGBColor(0xff, 0xf3, 0xe0)
LIGHT_RED = RGBColor(0xff, 0xeb, 0xee)
WHITE = RGBColor(0xff, 0xff, 0xff)

# =====================================================================
# Helpers (extended from generate_pptx.py)
# =====================================================================

def _text(box, text, size=18, bold=False, color=DARK, align=PP_ALIGN.LEFT, font_name=None):
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.alignment = align
    if font_name:
        p.font.name = font_name
    return tf


def add_text_box(slide, left, top, width, height, text,
                 font_size=18, bold=False, color=DARK,
                 align=PP_ALIGN.LEFT, font_name=None):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    _text(box, text, font_size, bold, color, align, font_name)
    return box


def add_bullet_list(slide, left, top, width, height, items,
                    font_size=16, color=DARK, spacing=6):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = f"•  {item}"
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.space_after = Pt(spacing)
    return box


def add_rect(slide, left, top, width, height,
             fill_color=None, line_color=None, shape=MSO_SHAPE.ROUNDED_RECTANGLE):
    sh = slide.shapes.add_shape(shape, Inches(left), Inches(top), Inches(width), Inches(height))
    if fill_color is not None:
        sh.fill.solid()
        sh.fill.fore_color.rgb = fill_color
    else:
        sh.fill.background()
    if line_color is not None:
        sh.line.color.rgb = line_color
    else:
        sh.line.fill.background()
    return sh


def add_title_slide(prs, title, subtitle=""):
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Top accent bar
    add_rect(slide, 0, 0, 13.333, 0.3, fill_color=BLUE, shape=MSO_SHAPE.RECTANGLE)
    # Bottom accent bar
    add_rect(slide, 0, 7.2, 13.333, 0.3, fill_color=DARK_BLUE, shape=MSO_SHAPE.RECTANGLE)

    add_text_box(slide, 0.5, 2.3, 12.333, 1.4, title, 54, bold=True, color=BLUE, align=PP_ALIGN.CENTER)
    if subtitle:
        add_text_box(slide, 0.5, 3.7, 12.333, 0.8, subtitle, 24, color=MUTED, align=PP_ALIGN.CENTER)

    return slide


def add_content_slide(prs, title, subtitle=""):
    """Create a slide with header bar + optional subtitle. Returns the slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Header bar
    add_rect(slide, 0, 0, 13.333, 1.0, fill_color=BLUE, shape=MSO_SHAPE.RECTANGLE)
    add_text_box(slide, 0.5, 0.22, 12.333, 0.6, title, 30, bold=True, color=WHITE)

    if subtitle:
        add_text_box(slide, 0.5, 1.15, 12.333, 0.4, subtitle, 16, color=MUTED)

    return slide


def add_stat_box(slide, left, top, number, label, color=BLUE):
    add_rect(slide, left, top, 2.2, 1.3, fill_color=color)
    add_text_box(slide, left, top + 0.15, 2.2, 0.6, number, 30, bold=True,
                 color=WHITE, align=PP_ALIGN.CENTER)
    add_text_box(slide, left, top + 0.75, 2.2, 0.45, label, 13,
                 color=WHITE, align=PP_ALIGN.CENTER)


def add_layer_box(slide, left, top, width, num, name, confidence, desc, color):
    add_rect(slide, left, top, width, 0.75, fill_color=color)
    add_text_box(slide, left + 0.2, top + 0.1, 0.6, 0.5, str(num),
                 24, bold=True, color=WHITE)
    add_text_box(slide, left + 0.8, top + 0.05, 5.0, 0.4, name,
                 16, bold=True, color=WHITE)
    add_text_box(slide, left + width - 2.7, top + 0.08, 2.5, 0.35, confidence,
                 12, color=WHITE, align=PP_ALIGN.RIGHT)
    add_text_box(slide, left + 0.8, top + 0.4, width - 1.0, 0.35, desc,
                 11, color=WHITE)


def add_chip(slide, left, top, text, fill=GREEN, width=None, height=0.35, size=11):
    if width is None:
        width = max(0.9, 0.11 * len(text) + 0.3)
    add_rect(slide, left, top, width, height, fill_color=fill)
    add_text_box(slide, left, top + (height - 0.28) / 2, width, 0.3, text,
                 size, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    return width


def add_arrow_down(slide, left, top, color=BLUE, width=0.5, height=0.55):
    sh = slide.shapes.add_shape(
        MSO_SHAPE.DOWN_ARROW,
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    sh.fill.solid()
    sh.fill.fore_color.rgb = color
    sh.line.fill.background()
    return sh


def add_arrow_right(slide, left, top, color=BLUE, width=0.65, height=0.45):
    sh = slide.shapes.add_shape(
        MSO_SHAPE.RIGHT_ARROW,
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    sh.fill.solid()
    sh.fill.fore_color.rgb = color
    sh.line.fill.background()
    return sh


# =====================================================================
# Build presentation
# =====================================================================

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)


# ---------------------------------------------------------------------
# SLIDE 1: Title
# ---------------------------------------------------------------------
slide = add_title_slide(
    prs,
    "LimeSurvey Data Pipeline",
    "From Raw Survey Chaos to Identity-Linked Insights"
)
add_stat_box(slide, 2.3, 5.0, "133", "Surveys")
add_stat_box(slide, 4.7, 5.0, "930K+", "Rows")
add_stat_box(slide, 7.1, 5.0, "251", "Standardized Questions", color=GREEN)
add_stat_box(slide, 9.5, 5.0, "~3,300", "Linkable Respondents", color=PURPLE)
add_text_box(slide, 0.5, 6.6, 12.333, 0.4,
             "BioNews Data Engineering  ·  Updated April 2026",
             14, color=MUTED, align=PP_ALIGN.CENTER)


# ---------------------------------------------------------------------
# SLIDE 2: Why This Matters
# ---------------------------------------------------------------------
slide = add_content_slide(prs, "Why This Matters",
                          "The business case in one slide")

# Left card: problem
add_rect(slide, 0.5, 1.7, 6.0, 5.4, fill_color=LIGHT_RED, line_color=RED)
add_chip(slide, 0.7, 1.88, "PROBLEM", fill=RED, width=1.1, height=0.38, size=10)
add_text_box(slide, 1.88, 1.85, 4.5, 0.5, "The Problem", 22, bold=True, color=RED)
add_bullet_list(slide, 0.7, 2.45, 5.7, 4.5, [
    "Same question asked 10 different ways across 133 surveys",
    "Answers spelled and capitalized inconsistently",
    "Impossible for a BI analyst to cleanly answer basic questions",
    "Respondents can't be matched to other data sources (Mailchimp, WordPress, GA4)",
    "~6,000 distinct question texts — far too many to review by hand",
], 15)

# Right card: solution
add_rect(slide, 6.8, 1.7, 6.0, 5.4, fill_color=LIGHT_GREEN, line_color=GREEN)
add_chip(slide, 7.0, 1.88, "SOLUTION", fill=GREEN, width=1.15, height=0.38, size=10)
add_text_box(slide, 8.23, 1.85, 4.5, 0.5, "The Solution", 22, bold=True, color=GREEN)
add_bullet_list(slide, 7.0, 2.45, 5.7, 4.5, [
    "Automated pipeline converts raw MySQL → analytics-ready BigQuery",
    "AI standardizes 251 canonical questions across all surveys",
    "Blended suggestion engine ranks matches from 4 signals",
    "Human-in-the-loop UI reviews edge cases with protection from overwrite",
    "Identity bridge links ~3,300 respondents to the full customer graph",
], 15)


# ---------------------------------------------------------------------
# SLIDE 3: 5-Stage Pipeline Overview
# ---------------------------------------------------------------------
slide = add_content_slide(prs, "The Pipeline in 5 Stages",
                          "Each stage has a clear purpose and a protected handoff")

stages = [
    ("1  EXTRACT",      "Pull raw survey data + tokens from LimeSurvey MySQL",            BLUE),
    ("2  ENRICH",       "Join question text & answer text — make it human-readable",       GREEN),
    ("3  STANDARDIZE",  "6-layer AI cascade + 4-signal suggestion engine",                 PURPLE),
    ("4  REVIEW",       "Human confirms edge cases via unified Streamlit mapper",          ORANGE),
    ("5  PUBLISH",      "Build wide analytics table + feed the Identity Hub",              RED),
]

top = 1.75
for title, desc, color in stages:
    add_rect(slide, 1.5, top, 10.3, 0.85, fill_color=color)
    add_text_box(slide, 1.7, top + 0.1, 3.2, 0.45, title, 20, bold=True, color=WHITE)
    add_text_box(slide, 5.0, top + 0.18, 6.8, 0.4, desc, 14, color=WHITE)
    top += 1.0


# ---------------------------------------------------------------------
# SLIDE 4: Data Flow Chart (the key slide)
# ---------------------------------------------------------------------
slide = add_content_slide(prs, "Data Flow Chart",
                          "Every data store + every process, end to end")

# Source (MySQL)
add_rect(slide, 0.3, 1.7, 2.3, 1.0, fill_color=GRAY)
add_text_box(slide, 0.3, 1.85, 2.3, 0.4, "LimeSurvey MySQL", 13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text_box(slide, 0.3, 2.25, 2.3, 0.4, "(source system)", 10, color=WHITE, align=PP_ALIGN.CENTER)

# Arrow
add_arrow_right(slide, 2.62, 2.05)

# Extractor
add_rect(slide, 3.2, 1.7, 2.3, 1.0, fill_color=BLUE)
add_text_box(slide, 3.2, 1.85, 2.3, 0.4, "Extractor", 13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text_box(slide, 3.2, 2.25, 2.3, 0.4, "SSH tunnel → Parquet", 10, color=WHITE, align=PP_ALIGN.CENTER)

# Arrow
add_arrow_right(slide, 5.52, 2.05)

# GCS
add_rect(slide, 6.1, 1.7, 2.3, 1.0, fill_color=GRAY)
add_text_box(slide, 6.1, 1.85, 2.3, 0.4, "GCS Staging", 13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text_box(slide, 6.1, 2.25, 2.3, 0.4, "parquet files", 10, color=WHITE, align=PP_ALIGN.CENTER)

# Arrow
add_arrow_right(slide, 8.42, 2.05)

# BQ Staging
add_rect(slide, 9.0, 1.7, 3.8, 1.0, fill_color=LIGHT_BLUE, line_color=BLUE)
add_text_box(slide, 9.2, 1.8, 3.4, 0.35, "lime_surveys_columnar", 13, bold=True, color=BLUE)
add_text_box(slide, 9.2, 2.15, 3.4, 0.3, "BQ STAGING · ~1M rows", 10, color=MUTED)
add_text_box(slide, 9.2, 2.4, 3.4, 0.3, "rebuildable any time", 10, color=MUTED)

# Arrow down
add_arrow_down(slide, 10.45, 2.72)

# Enrichment process
add_rect(slide, 8.5, 3.25, 4.3, 0.6, fill_color=GREEN)
add_text_box(slide, 8.5, 3.32, 4.3, 0.4, "Enrichment SQL (INSERT-only)", 13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# Arrow down
add_arrow_down(slide, 10.45, 3.87)

# BQ Production (protected) - centerpiece
add_rect(slide, 6.5, 4.4, 6.3, 1.3, fill_color=RED, line_color=DARK_BLUE)
add_chip(slide, 6.72, 4.52, "PROTECTED", fill=DARK_BLUE, width=1.3, height=0.36, size=9)
add_text_box(slide, 8.1, 4.5, 4.6, 0.45, "lime_surveys_columnar_completed", 15, bold=True, color=WHITE)
add_text_box(slide, 6.7, 4.95, 6.0, 0.3, "BQ PRODUCTION · SOURCE OF TRUTH · ~930K rows", 11, color=WHITE)
add_text_box(slide, 6.7, 5.25, 6.0, 0.3, "APPEND-ONLY · protected from --rebuild · manually_assigned lock", 11, color=WHITE)

# Tokens side-path
add_rect(slide, 0.3, 3.25, 2.3, 1.0, fill_color=GRAY)
add_text_box(slide, 0.3, 3.4, 2.3, 0.4, "lime_survey_tokens", 12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text_box(slide, 0.3, 3.8, 2.3, 0.4, "participant linkage", 10, color=WHITE, align=PP_ALIGN.CENTER)

# Processes running against completed
add_rect(slide, 0.3, 4.4, 2.9, 1.3, fill_color=PURPLE)
add_text_box(slide, 0.3, 4.55, 2.9, 0.35, "NLP Processor", 14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text_box(slide, 0.3, 4.9, 2.9, 0.3, "6-layer matching", 11, color=WHITE, align=PP_ALIGN.CENTER)
add_text_box(slide, 0.3, 5.2, 2.9, 0.3, "nightly retrain", 11, color=WHITE, align=PP_ALIGN.CENTER)

add_rect(slide, 3.3, 4.4, 3.0, 1.3, fill_color=ORANGE)
add_text_box(slide, 3.3, 4.55, 3.0, 0.35, "Streamlit Mapper", 14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text_box(slide, 3.3, 4.9, 3.0, 0.3, "4-signal suggestions", 11, color=WHITE, align=PP_ALIGN.CENTER)
add_text_box(slide, 3.3, 5.2, 3.0, 0.3, "auto-apply + undo", 11, color=WHITE, align=PP_ALIGN.CENTER)

# Arrow from completed down to wide
add_arrow_down(slide, 9.45, 5.72)

# Wide table + identity hub
add_rect(slide, 4.0, 6.25, 4.3, 1.0, fill_color=LIGHT_GREEN, line_color=GREEN)
add_text_box(slide, 4.15, 6.4, 4.0, 0.35, "lime_surveys_wide", 13, bold=True, color=GREEN)
add_text_box(slide, 4.15, 6.75, 4.0, 0.3, "ANALYTICS · ~63K rows · 1 per respondent", 10, color=MUTED)

add_arrow_right(slide, 8.42, 6.63, color=GREEN)

add_rect(slide, 9.0, 6.25, 3.8, 1.0, fill_color=PURPLE)
add_text_box(slide, 9.0, 6.4, 3.8, 0.35, "Identity Hub", 14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text_box(slide, 9.0, 6.75, 3.8, 0.3, "email + bnfpvid + profile fields", 10, color=WHITE, align=PP_ALIGN.CENTER)


# ---------------------------------------------------------------------
# SLIDE 5: Stage 1 — Extraction
# ---------------------------------------------------------------------
slide = add_content_slide(prs, "Stage 1: Extraction",
                          "Pull raw data from LimeSurvey MySQL into BigQuery")

add_text_box(slide, 0.5, 1.8, 6, 0.45, "What Happens", 22, bold=True, color=BLUE)
add_bullet_list(slide, 0.5, 2.35, 6, 3.5, [
    "Connects to LimeSurvey via secure SSH tunnel",
    "Parallel workers pull all 133 surveys at once",
    "Auto-discovers per-survey response tables",
    "Extracts token tables for participant linkage",
    "Runs nightly with a 7-day lookback window",
    "Arrives in 'long format' — one row per answer",
], 15)

add_text_box(slide, 7.0, 1.8, 5.8, 0.45, "Technologies", 22, bold=True, color=BLUE)
chips_top = 2.4
chip_left = 7.0
for c in ["Python", "SSH Tunnel", "MySQL", "BigQuery", "Parquet", "GCS"]:
    w = add_chip(slide, chip_left, chips_top, c, fill=GREEN, height=0.4, size=12)
    chip_left += w + 0.12
    if chip_left > 12.5:
        chip_left = 7.0
        chips_top += 0.55

# Output card
add_rect(slide, 7.0, 4.3, 5.8, 2.7, fill_color=LIGHT_BLUE, line_color=BLUE)
add_text_box(slide, 7.2, 4.45, 5.4, 0.4, "Output Tables", 16, bold=True, color=DARK)
add_text_box(slide, 7.2, 4.9, 5.4, 0.35, "lime_surveys_columnar", 14, bold=True, color=BLUE)
add_text_box(slide, 7.2, 5.25, 5.4, 0.3, "~1M rows · staging · rebuildable", 11, color=MUTED)
add_text_box(slide, 7.2, 5.7, 5.4, 0.35, "lime_survey_tokens", 14, bold=True, color=BLUE)
add_text_box(slide, 7.2, 6.05, 5.4, 0.3, "per-survey token → participant_id map", 11, color=MUTED)
add_text_box(slide, 7.2, 6.45, 5.4, 0.3, "(plus 12 metadata tables: questions, answers, users, etc.)", 11, color=MUTED)


# ---------------------------------------------------------------------
# SLIDE 6: Stage 2 — Enrichment
# ---------------------------------------------------------------------
slide = add_content_slide(prs, "Stage 2: Enrichment",
                          "Turn cryptic codes into human-readable text — safely")

# Before
add_rect(slide, 0.5, 1.8, 5.8, 2.0, fill_color=LIGHT_RED, line_color=RED)
add_text_box(slide, 0.7, 1.95, 5.5, 0.4, "Before (Raw)", 18, bold=True, color=RED)
add_text_box(slide, 0.7, 2.4, 5.5, 0.35, "question_id: 42817", 14, font_name="Courier New")
add_text_box(slide, 0.7, 2.75, 5.5, 0.35, "answer_value: AO01", 14, font_name="Courier New")
add_text_box(slide, 0.7, 3.15, 5.5, 0.3, "Not usable by analysts!", 12, color=RED)

add_arrow_right(slide, 6.35, 2.55, width=0.62, height=0.5)

# After
add_rect(slide, 7.0, 1.8, 5.8, 2.0, fill_color=LIGHT_GREEN, line_color=GREEN)
add_text_box(slide, 7.2, 1.95, 5.5, 0.4, "After (Enriched)", 18, bold=True, color=GREEN)
add_text_box(slide, 7.2, 2.4, 5.5, 0.35, 'raw_question_en: "What is your age?"', 14, font_name="Courier New")
add_text_box(slide, 7.2, 2.75, 5.5, 0.35, 'raw_answer_en: "18-24 years"', 14, font_name="Courier New")
add_text_box(slide, 7.2, 3.15, 5.5, 0.3, "Ready for standardization!", 12, color=GREEN)

# Key actions
add_rect(slide, 0.5, 4.1, 12.3, 2.7, fill_color=LIGHT_BLUE, line_color=BLUE)
add_text_box(slide, 0.7, 4.25, 11.9, 0.4, "What the Enrichment SQL Does", 18, bold=True, color=BLUE)
add_bullet_list(slide, 0.7, 4.75, 11.9, 2, [
    "JOINs question_id → readable question text (lime_question_l10ns)",
    "JOINs answer_value → readable answer text (lime_answers + lime_answer_l10ns)",
    "JOINs survey_id + token → participant_id (lime_survey_tokens)",
    "Strips HTML tags; leaves NULL for Stage 3 to fill in standardized_question/answer",
    "INSERT-only: never touches existing rows in lime_surveys_columnar_completed",
], 14)


# ---------------------------------------------------------------------
# SLIDE 7: Stage 3 — AI Standardization (6 layers)
# ---------------------------------------------------------------------
slide = add_content_slide(prs, "Stage 3: AI Standardization",
                          "A 6-layer cascade picks the best match — highest confidence wins")

layers = [
    (1, "EXACT MATCH",        "100% confidence",   "Already seen this exact question — reuse the prior mapping",           BLUE),
    (2, "MACHINE LEARNING",   "90–99% confidence", "Trained classifier predicts category from thousands of prior mappings", GREEN),
    (3, "SEMANTIC SIMILARITY","85–99% confidence", '"How old are you?" ≈ "What is your age?" via sentence embeddings',     ORANGE),
    (4, "MEDICAL TERMS (UMLS)","85–98% confidence",'"ALS" ≈ "Amyotrophic Lateral Sclerosis" via UMLS concept mapping',     RED),
    (5, "ENTITY RECOGNITION", "80–93% confidence", "spaCy identifies diseases, chemicals, body parts in the question",     PURPLE),
    (6, "FUZZY FALLBACK",     "70–85% confidence", "Catches typos and whitespace — last line of defense",                   GRAY),
]
top = 1.7
for num, name, conf, desc, color in layers:
    add_layer_box(slide, 0.5, top, 12.3, num, name, conf, desc, color)
    top += 0.87


# ---------------------------------------------------------------------
# SLIDE 8: The 4-Signal Blended Suggestion Engine
# ---------------------------------------------------------------------
slide = add_content_slide(prs, "4-Signal Blended Suggestion Engine",
                          "The Streamlit UI combines four independent signals to rank candidates")

signals = [
    ("1  FUZZY TEXT",        "RapidFuzz token-sort on question text",                        BLUE),
    ("2  SEMANTIC",          "Sentence-transformer embedding cosine similarity",              GREEN),
    ("3  CO-OCCURRENCE",     "Other questions sharing the same raw answer values",            ORANGE),
    ("4  ANSWER FINGERPRINT","Per-std-question answer profile (refreshed nightly)",           PURPLE),
]
top = 1.85
for title, desc, color in signals:
    add_rect(slide, 0.5, top, 12.3, 0.9, fill_color=color)
    add_text_box(slide, 0.7, top + 0.15, 3.5, 0.4, title, 18, bold=True, color=WHITE)
    add_text_box(slide, 4.3, top + 0.22, 8.3, 0.4, desc, 13, color=WHITE)
    top += 1.0

# Insight card
add_rect(slide, 0.5, 6.0, 12.3, 1.2, fill_color=LIGHT_AMBER, line_color=ORANGE)
add_text_box(slide, 0.7, 6.15, 11.9, 0.4, "Why 4 signals beat 1", 16, bold=True, color=ORANGE)
add_text_box(slide, 0.7, 6.55, 11.9, 0.55,
             "A candidate appearing in multiple signals outranks one appearing in just one. Example: "
             "\"Practitioner Details:\" → answer-set matches healthcare_specialists even though fuzzy text does not.",
             13, color=DARK)


# ---------------------------------------------------------------------
# SLIDE 9: Stage 4 — Unified Streamlit Mapper
# ---------------------------------------------------------------------
slide = add_content_slide(prs, "Stage 4: Human Review",
                          "Unified Streamlit mapper — table + detail, one question at a time")

# Left: Features
add_text_box(slide, 0.5, 1.85, 6, 0.45, "Key Features", 20, bold=True, color=BLUE)
add_bullet_list(slide, 0.5, 2.4, 6, 4.7, [
    "Table on the left ranks the queue by suggestion score",
    "Detail panel on the right shows full context + candidate list",
    "One-click Accept for high-confidence matches (≥95%)",
    "Bulk auto-accept at user-set threshold",
    "Near-duplicate clustering — map 10 variants in one click",
    "Inline category creation — never leave the screen",
    "Live streaming logs for every background job",
    "Keyboard-friendly workflow",
], 14)

# Right: Benefits
add_rect(slide, 7.0, 1.85, 5.8, 5.3, fill_color=LIGHT_GREEN, line_color=GREEN)
add_text_box(slide, 7.2, 2.0, 5.5, 0.4, "Benefits", 20, bold=True, color=GREEN)
add_bullet_list(slide, 7.2, 2.55, 5.5, 4.5, [
    "Reviewer throughput: dozens → hundreds per hour",
    "Manual assignments protected from ETL overwrite",
    "Full undo + undo-batch — reviewers can be bold",
    "Audit log of every decision with timestamp",
    "Replaces 3 overlapping apps with one consistent workflow",
    "Suspicious-mapping warnings before commit",
], 14)


# ---------------------------------------------------------------------
# SLIDE 10: Stage 5 — Analytics Output
# ---------------------------------------------------------------------
slide = add_content_slide(prs, "Stage 5: Analytics Output",
                          "Long format → Wide format for BI + Identity Hub")

# Before
add_text_box(slide, 0.5, 1.85, 6, 0.45, "Long Format (completed)", 16, bold=True, color=DARK)
add_text_box(slide, 0.5, 2.3, 6, 0.3, "One row per answer  ·  ~930K rows", 12, color=MUTED)

# Mini table
headers = ["response_id", "std_question", "std_answer"]
rows = [["12345", "age", "25_to_34"], ["12345", "gender", "female"], ["12345", "diagnosis", "als"]]
t = 2.7
for i, h in enumerate(headers):
    add_rect(slide, 0.5 + i * 1.9, t, 1.85, 0.4, fill_color=BLUE, shape=MSO_SHAPE.RECTANGLE)
    add_text_box(slide, 0.55 + i * 1.9, t + 0.06, 1.75, 0.3, h, 11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
t += 0.45
for r in rows:
    for i, c in enumerate(r):
        add_text_box(slide, 0.55 + i * 1.9, t, 1.75, 0.32, c, 11, align=PP_ALIGN.CENTER, font_name="Courier New")
    t += 0.35

add_arrow_right(slide, 6.42, 3.35, width=0.62, height=0.5)

# After
add_text_box(slide, 7.0, 1.85, 5.8, 0.45, "Wide Format (lime_surveys_wide)", 16, bold=True, color=DARK)
add_text_box(slide, 7.0, 2.3, 5.8, 0.3, "One row per respondent  ·  ~63K rows  ·  ~150 columns", 12, color=MUTED)

headers2 = ["response_id", "age", "gender", "diagnosis"]
rows2 = [["12345", "25_to_34", "female", "als"], ["12346", "45_to_54", "male", "ms"]]
t = 2.7
for i, h in enumerate(headers2):
    add_rect(slide, 7.0 + i * 1.45, t, 1.4, 0.4, fill_color=BLUE, shape=MSO_SHAPE.RECTANGLE)
    add_text_box(slide, 7.05 + i * 1.45, t + 0.06, 1.3, 0.3, h, 11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
t += 0.45
for r in rows2:
    for i, c in enumerate(r):
        add_text_box(slide, 7.05 + i * 1.45, t, 1.3, 0.32, c, 11, align=PP_ALIGN.CENTER, font_name="Courier New")
    t += 0.35

# Downstream uses
add_rect(slide, 0.5, 4.75, 12.3, 2.4, fill_color=LIGHT_BLUE, line_color=BLUE)
add_text_box(slide, 0.7, 4.9, 11.9, 0.4, "Who consumes lime_surveys_wide?", 17, bold=True, color=BLUE)
add_bullet_list(slide, 0.7, 5.4, 11.9, 1.7, [
    "BI analysts + Looker dashboards — straightforward SQL with one column per question",
    "Profile enrichment (Parts 11 & 13) — zero-party clinical data fills 0–1%-fill fields",
    "Identity Hub — email bridge to Mailchimp/WordPress, GA4 session matching",
], 13)


# ---------------------------------------------------------------------
# SLIDE 11: Identity Hub Integration
# ---------------------------------------------------------------------
slide = add_content_slide(prs, "Identity Hub Integration",
                          "Survey respondents become the highest-quality per-profile data source")

add_text_box(slide, 0.5, 1.85, 12.3, 0.4,
             "The bridge is self-reported email in survey answers, not the lime_participants table.",
             14, color=MUTED)

# Path diagram
add_rect(slide, 0.5, 2.4, 3.2, 1.2, fill_color=BLUE)
add_text_box(slide, 0.5, 2.55, 3.2, 0.4, "Survey Response", 14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text_box(slide, 0.5, 2.95, 3.2, 0.4, "(survey_id, response_id)", 11, color=WHITE, align=PP_ALIGN.CENTER)
add_text_box(slide, 0.5, 3.25, 3.2, 0.3, "standardized_question='email'", 10, color=WHITE, align=PP_ALIGN.CENTER)

add_arrow_right(slide, 3.77, 2.83)

add_rect(slide, 4.4, 2.4, 3.2, 1.2, fill_color=GREEN)
add_text_box(slide, 4.4, 2.55, 3.2, 0.4, "Email Bridge", 14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text_box(slide, 4.4, 2.95, 3.2, 0.4, "Phase 2", 11, color=WHITE, align=PP_ALIGN.CENTER)
add_text_box(slide, 4.4, 3.25, 3.2, 0.3, "Mailchimp, WordPress", 10, color=WHITE, align=PP_ALIGN.CENTER)

add_arrow_right(slide, 7.67, 2.83)

add_rect(slide, 8.3, 2.4, 4.4, 1.2, fill_color=PURPLE)
add_text_box(slide, 8.3, 2.55, 4.4, 0.4, "Session Connector", 14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text_box(slide, 8.3, 2.95, 4.4, 0.4, "Phase 5 — 5-min GA4 match", 11, color=WHITE, align=PP_ALIGN.CENTER)
add_text_box(slide, 8.3, 3.25, 4.4, 0.3, "bnfpvid, client_id", 10, color=WHITE, align=PP_ALIGN.CENTER)

# Stats card
add_rect(slide, 0.5, 4.0, 12.3, 1.4, fill_color=LIGHT_GREEN, line_color=GREEN)
add_text_box(slide, 0.7, 4.15, 11.9, 0.4, "What we have today", 16, bold=True, color=GREEN)
add_text_box(slide, 0.7, 4.6, 11.9, 0.4, "4,548 respondents with self-reported email  ·  4,065 with last_name  ·  2,775 with first_name", 13, color=DARK)
add_text_box(slide, 0.7, 4.95, 11.9, 0.4, "~3,300 linkable to Mailchimp / WordPress / GA4  ·  growing every day", 13, color=DARK)

# Value card
add_rect(slide, 0.5, 5.6, 12.3, 1.5, fill_color=LIGHT_AMBER, line_color=ORANGE)
add_text_box(slide, 0.7, 5.75, 11.9, 0.4, "Why survey data is the highest-value source", 16, bold=True, color=ORANGE)
add_bullet_list(slide, 0.7, 6.2, 11.9, 0.9, [
    "Zero-party data — respondents explicitly self-identify condition, treatments, diagnosis timing",
    "Fills profile fields that sit at 0–1% fill rate from every other source",
], 12)


# ---------------------------------------------------------------------
# SLIDE 12: Data Protection
# ---------------------------------------------------------------------
slide = add_content_slide(prs, "Data Protection",
                          "Three independent safety layers protect your work")

protections = [
    ("PROTECTED_TABLES list",   "lime_surveys_columnar_completed is hardcoded as protected. --rebuild cannot drop it.", RED),
    ("manually_assigned flag",  "Every manual mapping sets this flag. The ETL refuses to overwrite flagged rows.",       ORANGE),
    ("mapping_history audit",   "Every save is logged with before/after values. Full Undo + Undo-batch available.",      GREEN),
    ("Append-only INSERTs",     "Enrichment SQL only inserts rows with submitdate > MAX(existing). Never deletes.",       BLUE),
]
top = 1.85
for title, desc, color in protections:
    add_rect(slide, 0.5, top, 12.3, 1.15, fill_color=color)
    add_text_box(slide, 0.7, top + 0.15, 11.9, 0.45, title, 18, bold=True, color=WHITE)
    add_text_box(slide, 0.7, top + 0.6, 11.9, 0.5, desc, 13, color=WHITE)
    top += 1.25

# Footer
add_text_box(slide, 0.5, 6.85, 12.3, 0.4,
             "Net effect: reviewers can experiment, try bulk actions, and recover — nothing is permanently destructive.",
             12, bold=True, color=MUTED, align=PP_ALIGN.CENTER)


# ---------------------------------------------------------------------
# SLIDE 13: Daily Operations
# ---------------------------------------------------------------------
slide = add_content_slide(prs, "Daily Operations",
                          "What runs automatically vs. what needs human attention")

# Auto column
add_rect(slide, 0.5, 1.85, 6.0, 5.3, fill_color=LIGHT_BLUE, line_color=BLUE)
add_chip(slide, 0.7, 2.0, "AUTO", fill=BLUE, width=0.7, height=0.38, size=10)
add_text_box(slide, 1.48, 1.98, 5.0, 0.42, "Runs Automatically (Nightly)", 17, bold=True, color=BLUE)
add_bullet_list(slide, 0.7, 2.55, 5.7, 4.5, [
    "Extract new responses (7-day lookback)",
    "Extract token tables",
    "Enrichment SQL INSERT into completed",
    "6-layer NLP standardization on NULL rows",
    "ML model retrain if training pool grew 10%+",
    "Answer fingerprint rebuild",
    "Wide table full rebuild",
    "Coverage stats + health check",
], 13)

# Manual column
add_rect(slide, 6.8, 1.85, 6.0, 5.3, fill_color=LIGHT_AMBER, line_color=ORANGE)
add_chip(slide, 7.0, 2.0, "MANUAL", fill=ORANGE, width=0.85, height=0.38, size=10)
add_text_box(slide, 7.93, 1.98, 4.8, 0.42, "Needs Human Attention", 17, bold=True, color=ORANGE)
add_bullet_list(slide, 7.0, 2.55, 5.7, 4.5, [
    "Review suggestions the AI couldn't resolve",
    "Accept bulk high-confidence matches",
    "Add new standardized questions as surveys evolve",
    "Review low-confidence bands periodically",
    "Create new categories when needed",
    "Respond to suspicious-mapping warnings",
    "Undo mistakes (always available)",
], 13)


# ---------------------------------------------------------------------
# SLIDE 14: Summary
# ---------------------------------------------------------------------
slide = add_content_slide(prs, "Summary",
                          "One pipeline, many payoffs")

# Top row: key metrics
add_stat_box(slide, 1.0, 1.9, "99.96%", "Questions Standardized", color=GREEN)
add_stat_box(slide, 3.5, 1.9, "930K+", "Protected Rows", color=BLUE)
add_stat_box(slide, 6.0, 1.9, "251", "Canonical Questions", color=PURPLE)
add_stat_box(slide, 8.5, 1.9, "~3.3K", "Linkable Profiles", color=ORANGE)
add_stat_box(slide, 11.0, 1.9, "4", "Suggestion Signals", color=RED)

# What you get
add_rect(slide, 0.5, 3.6, 12.3, 3.2, fill_color=LIGHT_GREEN, line_color=GREEN)
add_text_box(slide, 0.7, 3.75, 11.9, 0.45, "What you get", 20, bold=True, color=GREEN)
add_bullet_list(slide, 0.7, 4.3, 11.9, 2.4, [
    "BI-ready wide table — one row per respondent, ~150 canonical columns",
    "Deterministic identity bridge linking ~3,300 respondents to the full customer graph",
    "Zero-party clinical data (condition, treatment, diagnosis timing) at far higher fill rate than any other source",
    "Reviewer workflow that scales: bulk auto-accept + undo + keyboard-fast detail pane",
    "Belt-and-suspenders data protection — nothing destructive can happen accidentally",
], 14)

# Footer
add_text_box(slide, 0.5, 7.0, 12.3, 0.4,
             "The source keeps improving as reviewers work through the queue — quality compounds daily.",
             12, bold=True, color=MUTED, align=PP_ALIGN.CENTER)


# =====================================================================
# Save
# =====================================================================
out_path = Path(__file__).parent / "LimeSurvey_Pipeline_v2.pptx"
prs.save(out_path)
print(f"[OK] Saved presentation: {out_path}")
print(f"     Slides: {len(prs.slides)}")
