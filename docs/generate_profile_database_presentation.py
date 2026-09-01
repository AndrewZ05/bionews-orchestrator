"""
Generate the BioNews Profile Database PowerPoint presentation.

Matches the visual style of docs/Identity_Hub_Pipeline_Presentation_v5.pptx:
  - 16:9 (13.33" x 7.5")
  - Background F8F9FA, header bar 0D47A1 (1.35" tall) + accent stripe 1A73E8 (0.06")
  - Body text 212121, muted 607D8B, accent dark 0D2540, white cards FFFFFF
  - Title slide uses 0D47A1 background with deeper 0D2540 band

Use-case navigator slide has 12 clickable cards arranged 2 columns x 6 rows;
each card hyperlinks via shape.click_action.target_slide to a per-use-case SQL
slide. Each SQL slide has a "All Examples" back-link.

Run:
    python docs/generate_profile_database_presentation.py
Output:
    docs/Profile_Database_Presentation.pptx
"""
from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


# --------------------------------------------------------------------------
# Theme
# --------------------------------------------------------------------------
PRIMARY = RGBColor(0x0D, 0x47, 0xA1)        # deep blue
ACCENT = RGBColor(0x1A, 0x73, 0xE8)         # bright blue
DEEP = RGBColor(0x0D, 0x25, 0x40)           # almost-black blue (cover band)
BG = RGBColor(0xF8, 0xF9, 0xFA)             # off-white
BODY = RGBColor(0x21, 0x21, 0x21)           # body text
MUTED = RGBColor(0x60, 0x7D, 0x8B)          # muted text
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
CARD = RGBColor(0xFF, 0xFF, 0xFF)
BORDER = RGBColor(0xE0, 0xE0, 0xE0)         # subtle card border
GREEN = RGBColor(0x2E, 0x7D, 0x32)
ORANGE = RGBColor(0xEF, 0x6C, 0x00)
CODE_BG = RGBColor(0x0D, 0x25, 0x40)        # dark code-block background
CODE_TXT = RGBColor(0xE3, 0xF2, 0xFD)       # very light blue text


SLIDE_W = Inches(13.3333)
SLIDE_H = Inches(7.5)
HEADER_H = Inches(1.35)
STRIPE_H = Inches(0.06)


def set_solid(shape, rgb: RGBColor) -> None:
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb


def no_outline(shape) -> None:
    shape.line.fill.background()


def add_rect(slide, x, y, w, h, fill_rgb, *, outline=None):
    sh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    set_solid(sh, fill_rgb)
    if outline is None:
        no_outline(sh)
    else:
        sh.line.color.rgb = outline
        sh.line.width = Pt(0.75)
    return sh


def add_text(
    slide, x, y, w, h, text, *, size=14, bold=False, color=BODY,
    align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, name="Body",
):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tb.name = name
    tf = tb.text_frame
    tf.margin_left = Pt(0)
    tf.margin_right = Pt(0)
    tf.margin_top = Pt(0)
    tf.margin_bottom = Pt(0)
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    lines = str(text).split("\n")
    for i, line in enumerate(lines):
        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        para.alignment = align
        run = para.add_run()
        run.text = line
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color
        run.font.name = "Calibri"
    return tb


def add_header(slide, title, subtitle=""):
    """Standard slide header bar 1.35" tall + 0.06" accent stripe."""
    add_rect(slide, 0, 0, SLIDE_W, HEADER_H, PRIMARY)
    add_text(
        slide, Inches(0.40), Inches(0.12), Inches(12.5), Inches(0.65),
        title, size=28, bold=True, color=WHITE,
    )
    if subtitle:
        add_text(
            slide, Inches(0.40), Inches(0.78), Inches(12.5), Inches(0.45),
            subtitle, size=14, color=WHITE,
        )
    add_rect(slide, 0, HEADER_H, SLIDE_W, STRIPE_H, ACCENT)


def add_footer(slide, text):
    add_text(
        slide, Inches(0.40), Inches(7.05), Inches(12.5), Inches(0.40),
        text, size=10, color=MUTED, anchor=MSO_ANCHOR.MIDDLE,
    )


# --------------------------------------------------------------------------
# Slide builders
# --------------------------------------------------------------------------
def slide_cover(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_rect(s, 0, 0, SLIDE_W, SLIDE_H, PRIMARY)
    add_rect(s, 0, Inches(5.6), SLIDE_W, Inches(1.9), DEEP)
    add_rect(s, 0, Inches(5.6), SLIDE_W, Inches(0.06), ACCENT)
    add_text(
        s, Inches(0.6), Inches(2.1), Inches(12.0), Inches(1.0),
        "BioNews", size=40, bold=True, color=WHITE,
    )
    add_text(
        s, Inches(0.6), Inches(3.0), Inches(12.0), Inches(1.2),
        "Profile Database", size=56, bold=True, color=WHITE,
    )
    add_text(
        s, Inches(0.6), Inches(4.2), Inches(12.0), Inches(0.6),
        "v2.2 — Persona-resolved customer profiles built on the identity graph",
        size=18, color=WHITE,
    )
    add_text(
        s, Inches(0.6), Inches(5.85), Inches(12.0), Inches(0.5),
        "Architecture, Data Flows, Use Cases & SQL Reference",
        size=18, color=WHITE,
    )
    add_text(
        s, Inches(0.6), Inches(6.4), Inches(12.0), Inches(0.4),
        "Version 2.2  |  April 2026", size=14, color=MUTED,
    )
    add_text(
        s, Inches(0.6), Inches(7.05), Inches(12.0), Inches(0.4),
        "CONFIDENTIAL — Bionews Internal Use Only",
        size=10, color=MUTED,
    )
    return s


def slide_agenda(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_rect(s, 0, 0, SLIDE_W, SLIDE_H, BG)
    add_header(s, "Agenda", "What this deck covers")
    items = [
        ("01", "What Is the Profile Database?",
         "Resolved customer profiles on top of the identity graph"),
        ("02", "Production Snapshot",
         "3.46M profiles, 18M identifiers, 6.4M ad-attribution rows"),
        ("03", "Architecture & Data Flows",
         "Sources → populate → fill_gaps → enrich → personas"),
        ("04", "Output Tables",
         "16 consumer tables, 7 ops tables, 22 governed views"),
        ("05", "Build Modes",
         "rebuild • resume_rebuild • refresh • reenrich • views"),
        ("06", "Blue/Green Release",
         "Candidate publish, gated promotion, release manifest"),
        ("07", "Quality Gates & Observability",
         "Hard + soft assertions, runtime fingerprint, performance checks"),
        ("08", "Audience Surfaces",
         "Patient / HCP / Caregiver / Marketing / Analytics views"),
        ("09", "Features & Benefits",
         "What downstream consumers get out of the box"),
        ("10", "Use Cases with SQL",
         "Click-through reference for the most common queries"),
        ("11", "Operations Runbook",
         "Daily refresh, full rebuild, recovery from partial publish"),
    ]
    for i, (num, title, desc) in enumerate(items):
        col = i // 6
        row = i % 6
        x = Inches(0.28 + col * 6.66)
        y = Inches(1.55 + row * 0.90)
        add_rect(s, x, y, Inches(6.10), Inches(0.84), CARD, outline=BORDER)
        add_text(
            s, x + Inches(0.12), y + Inches(0.10),
            Inches(0.50), Inches(0.56),
            num, size=18, bold=True, color=ACCENT, anchor=MSO_ANCHOR.MIDDLE,
        )
        add_text(
            s, x + Inches(0.60), y + Inches(0.10),
            Inches(5.40), Inches(0.32),
            title, size=12, bold=True, color=PRIMARY,
        )
        add_text(
            s, x + Inches(0.60), y + Inches(0.46),
            Inches(5.40), Inches(0.30),
            desc, size=9, color=MUTED,
        )
    return s


def slide_what_is(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_rect(s, 0, 0, SLIDE_W, SLIDE_H, BG)
    add_header(
        s, "What Is the Profile Database?",
        "A persona-resolved customer-360 layer built on top of the identity graph",
    )
    add_text(
        s, Inches(0.40), Inches(1.55), Inches(12.5), Inches(0.50),
        "The identity graph stitches identifiers; the profile database resolves them into actionable customer profiles.",
        size=13, color=BODY,
    )
    boxes = [
        ("Identity Graph", "Stitches 16 identifier types into stable bn_ids",
         "identity_hub_data.bn_id_xref"),
        ("Profile Database", "Adds persona, condition, engagement, consent, signals",
         "profile_data.profile_core + 22 views"),
        ("Activation", "Audiences, contactability, marketing & analytics surfaces",
         "profile_marketing_audience, profile_audience_*"),
    ]
    for i, (title, desc, footer) in enumerate(boxes):
        x = Inches(0.40 + i * 4.30)
        y = Inches(2.30)
        w = Inches(4.10)
        add_rect(s, x, y, w, Inches(2.5), CARD, outline=BORDER)
        add_rect(s, x, y, w, Inches(0.55), PRIMARY)
        add_text(
            s, x + Inches(0.20), y + Inches(0.10), w - Inches(0.40), Inches(0.40),
            title, size=15, bold=True, color=WHITE,
        )
        add_text(
            s, x + Inches(0.20), y + Inches(0.75), w - Inches(0.40), Inches(0.90),
            desc, size=12, color=BODY,
        )
        add_text(
            s, x + Inches(0.20), y + Inches(1.85), w - Inches(0.40), Inches(0.50),
            footer, size=10, color=ACCENT, name="Code",
        )

    add_text(
        s, Inches(0.40), Inches(5.10), Inches(12.5), Inches(0.40),
        "Key principles", size=14, bold=True, color=PRIMARY,
    )
    bullets = [
        "• SQL-first: pure BigQuery; no app code in the resolution path",
        "• Immutable input contract: identity graph is upstream and not mutated",
        "• Persona-aware: every profile carries account_type, persona_source, and confidence",
        "• Governed surface: profile_current_safe is the default; profile_current is approved-only",
        "• Observable: every build emits run + step rows in profile_ops with runtime fingerprint",
    ]
    for i, b in enumerate(bullets):
        add_text(
            s, Inches(0.40), Inches(5.45 + i * 0.30), Inches(12.5), Inches(0.30),
            b, size=12, color=BODY,
        )
    add_footer(s, "Default consumer surface: profile_data.profile_current_safe")
    return s


def slide_scale(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_rect(s, 0, 0, SLIDE_W, SLIDE_H, BG)
    add_header(
        s, "Production Snapshot",
        "As of April 2026 — production dataset bi-data-391216.profile_data",
    )
    metrics = [
        ("3,460,572", "Total profiles", "Rows in profile_core"),
        ("411,916", "Marketable", "Email + opt-in"),
        ("305,941", "HCP audience", "NPI-anchored + verified"),
        ("17,649", "Confirmed patients", "account_type=patient"),
        ("18,161,373", "Identifier rows", "profile_identifiers"),
        ("6,437,821", "Ad-attribution rows", "profile_ad_attribution"),
        ("2,663,842", "Meta Pixel reach", "Facebook Custom Audience eligible"),
        ("2,590,103", "Google Ads reach", "Customer Match eligible"),
        ("2,352,718", "Content affinity rows", "profile_content_affinity"),
        ("2,619,972", "Have preferred condition", "75.7% of profiles"),
        ("3,464,071", "Engagement rows", "profile_engagement"),
        ("22", "Governed views", "profile_current_safe + audience surfaces"),
    ]
    for i, (val, label, sub) in enumerate(metrics):
        col = i % 4
        row = i // 4
        x = Inches(0.30 + col * 3.20)
        y = Inches(1.65 + row * 1.65)
        add_rect(s, x, y, Inches(3.00), Inches(1.45), CARD, outline=BORDER)
        add_text(
            s, x + Inches(0.10), y + Inches(0.12),
            Inches(2.80), Inches(0.55),
            val, size=24, bold=True, color=PRIMARY,
        )
        add_text(
            s, x + Inches(0.10), y + Inches(0.70),
            Inches(2.80), Inches(0.30),
            label, size=12, bold=True, color=BODY,
        )
        add_text(
            s, x + Inches(0.10), y + Inches(1.02),
            Inches(2.80), Inches(0.32),
            sub, size=10, color=MUTED,
        )
    add_footer(
        s,
        "Snapshot from build 7a94edbc (2026-04-25). Reproduce: SELECT COUNT(*) FROM profile_data.<table>",
    )
    return s


def slide_persona_mix(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_rect(s, 0, 0, SLIDE_W, SLIDE_H, BG)
    add_header(
        s, "Persona Mix & Engagement Distribution",
        "Where the 3.46M profiles sit in the persona / engagement matrix",
    )
    # Persona mix chart bars
    add_text(
        s, Inches(0.40), Inches(1.65), Inches(6.0), Inches(0.40),
        "Account type", size=14, bold=True, color=PRIMARY,
    )
    persona = [
        ("other / anonymous-known", 3100190, 89.6),
        ("hcp", 340914, 9.85),
        ("patient", 17649, 0.51),
        ("caregiver", 1241, 0.04),
        ("family_or_friend", 578, 0.02),
    ]
    for i, (label, n, pct) in enumerate(persona):
        y = Inches(2.10 + i * 0.55)
        add_text(
            s, Inches(0.40), y, Inches(2.4), Inches(0.30),
            label, size=11, color=BODY,
        )
        bar_w = max(0.05, min(3.4, pct / 100 * 3.4))
        add_rect(s, Inches(2.85), y + Inches(0.05), Inches(bar_w), Inches(0.22), ACCENT)
        add_text(
            s, Inches(6.30), y, Inches(0.6), Inches(0.30),
            f"{pct}%", size=11, bold=True, color=PRIMARY,
        )

    # Engagement tier
    add_text(
        s, Inches(7.10), Inches(1.65), Inches(6.0), Inches(0.40),
        "Engagement tier", size=14, bold=True, color=PRIMARY,
    )
    tiers = [
        ("low", 2428406, 70.1, MUTED),
        ("inactive", 595798, 17.2, RGBColor(0x9E, 0x9E, 0x9E)),
        ("medium", 246701, 7.1, ORANGE),
        ("high", 193166, 5.6, GREEN),
    ]
    for i, (label, n, pct, color) in enumerate(tiers):
        y = Inches(2.10 + i * 0.55)
        add_text(
            s, Inches(7.10), y, Inches(2.0), Inches(0.30),
            label, size=11, color=BODY,
        )
        bar_w = max(0.05, min(3.4, pct / 100 * 3.4))
        add_rect(s, Inches(9.10), y + Inches(0.05), Inches(bar_w), Inches(0.22), color)
        add_text(
            s, Inches(12.55), y, Inches(0.6), Inches(0.30),
            f"{pct}%", size=11, bold=True, color=PRIMARY,
        )

    # Top conditions
    add_text(
        s, Inches(0.40), Inches(5.05), Inches(12.5), Inches(0.40),
        "Top preferred conditions (74% of profiles carry one)",
        size=14, bold=True, color=PRIMARY,
    )
    conds = [
        ("MS", "404,688"),
        ("Myasthenia G.", "238,282"),
        ("Parkinson's", "184,635"),
        ("ALS", "178,217"),
        ("Hemophilia", "174,847"),
        ("Auto-Encep.", "92,968"),
        ("Amyloidosis", "92,365"),
        ("Renal Cell C.", "79,547"),
        ("MD", "75,991"),
        ("SMA", "73,017"),
    ]
    for i, (label, val) in enumerate(conds):
        col = i % 5
        row = i // 5
        x = Inches(0.40 + col * 2.55)
        y = Inches(5.55 + row * 0.75)
        add_rect(s, x, y, Inches(2.40), Inches(0.65), CARD, outline=BORDER)
        add_text(
            s, x + Inches(0.10), y + Inches(0.05), Inches(2.20), Inches(0.30),
            label, size=11, bold=True, color=BODY,
        )
        add_text(
            s, x + Inches(0.10), y + Inches(0.32), Inches(2.20), Inches(0.30),
            val, size=14, bold=True, color=ACCENT,
        )
    return s


def slide_architecture(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_rect(s, 0, 0, SLIDE_W, SLIDE_H, BG)
    add_header(
        s, "Architecture & Data Flow",
        "Five-stage SQL pipeline producing a governed customer-360 surface",
    )
    stages = [
        ("Sources", "Identity hub, Mailchimp, WordPress,\nLimeSurvey, GA4, NPI registry,\nOneTrust, AIM clickstream",
         RGBColor(0x42, 0xA5, 0xF5)),
        ("Populate", "Identity core CTAS\nEngagement summary\nSurvey, prefs, finalize",
         ACCENT),
        ("Fill Gaps", "Condition + site domain\nGA4 + forms + tier\nAIM + ads + affinity",
         RGBColor(0x29, 0xB6, 0xF6)),
        ("Enrich", "Subtypes / diagnosis\nMailchimp + GA4 lifetime\nHCP + zero-party",
         PRIMARY),
        ("Personas", "account_type + condition\npatient / caregiver detail\nstage + completeness",
         DEEP),
    ]
    box_w = Inches(2.40)
    box_h = Inches(2.20)
    gap = Inches(0.10)
    arrow_w = Inches(0.20)
    total = box_w * 5 + gap * 4 + arrow_w * 4
    start_x = (SLIDE_W - total) / 2
    y = Inches(2.0)
    for i, (title, body, color) in enumerate(stages):
        x = start_x + i * (box_w + gap + arrow_w)
        add_rect(s, x, y, box_w, box_h, CARD, outline=BORDER)
        add_rect(s, x, y, box_w, Inches(0.50), color)
        add_text(
            s, x + Inches(0.15), y + Inches(0.08), box_w - Inches(0.30), Inches(0.40),
            title, size=14, bold=True, color=WHITE,
        )
        add_text(
            s, x + Inches(0.15), y + Inches(0.65), box_w - Inches(0.30), Inches(1.50),
            body, size=10, color=BODY,
        )
        if i < 4:
            ax = x + box_w + Inches(0.05)
            ay = y + box_h / 2 - Inches(0.10)
            arrow = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, ax, ay, Inches(0.20), Inches(0.20))
            set_solid(arrow, MUTED)
            no_outline(arrow)

    add_text(
        s, Inches(0.40), Inches(4.50), Inches(12.5), Inches(0.45),
        "After personas: restore (app-authored fields), build candidate views, run gates, then promote.",
        size=12, color=BODY,
    )
    add_text(
        s, Inches(0.40), Inches(5.10), Inches(12.5), Inches(0.40),
        "Blue/green for rebuild", size=14, bold=True, color=PRIMARY,
    )
    add_text(
        s, Inches(0.40), Inches(5.50), Inches(12.5), Inches(1.6),
        "1.  Stage candidate tables in profile_data_candidate\n"
        "2.  Build candidate views in profile_staging\n"
        "3.  Run post-build assertions (6 hard, 5 soft) against the candidate\n"
        "4.  Repoint production views to candidate dataset\n"
        "5.  Copy 16 physical tables candidate → profile_data (logged in profile_publish_manifest)\n"
        "6.  Repoint production views back to profile_data and write profile_core_snapshot",
        size=12, color=BODY,
    )
    return s


def slide_output_tables(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_rect(s, 0, 0, SLIDE_W, SLIDE_H, BG)
    add_header(
        s, "Output Tables & Views",
        "Three tiers: consumer surface, derived signals, internal runtime",
    )
    columns = [
        ("Consumer Core", "profile_data", PRIMARY, [
            ("profile_core", "Persona, condition, demographics, HCP fields"),
            ("profile_identifiers", "All known IDs per profile (cross-system join)"),
            ("profile_engagement", "Email, web, BuddyPress, ad-click rollup"),
            ("profile_preferences", "Newsletter + forum settings"),
            ("profile_survey_data", "LimeSurvey S5 normalized answers"),
            ("site_events", "GA4-backed site interactions"),
            ("profile_zero_party", "Identity-linked poll/quiz answers"),
        ]),
        ("Derived Signals", "profile_data", ACCENT, [
            ("profile_content_affinity", "Per-condition browsing affinity score"),
            ("profile_ad_attribution", "Click-ID + platform + first/last-seen"),
            ("profile_segment_tags", "Governed Mailchimp + manual tags"),
            ("conditions_dict + 3 more", "Curated MeSH / RxNorm / SNOMED dicts"),
        ]),
        ("Ops & Staging", "profile_ops + profile_staging", DEEP, [
            ("profile_build_runs", "Run-level observability"),
            ("profile_build_steps", "Step-level performance metrics"),
            ("profile_publish_manifest", "Per-table promotion records"),
            ("profile_core_snapshot", "Point-in-time profile_core copy"),
            ("profile_field_changes", "Narrow audit log (4 persona fields)"),
            ("profile_restore_unmapped", "Snapshot rows that did not remap"),
        ]),
    ]
    col_w = Inches(4.20)
    col_h = Inches(5.40)
    gap = Inches(0.20)
    start_x = (SLIDE_W - col_w * 3 - gap * 2) / 2
    y = Inches(1.60)
    for i, (title, dataset, color, items) in enumerate(columns):
        x = start_x + i * (col_w + gap)
        add_rect(s, x, y, col_w, col_h, CARD, outline=BORDER)
        add_rect(s, x, y, col_w, Inches(0.65), color)
        add_text(
            s, x + Inches(0.15), y + Inches(0.05), col_w - Inches(0.30), Inches(0.30),
            title, size=14, bold=True, color=WHITE,
        )
        add_text(
            s, x + Inches(0.15), y + Inches(0.36), col_w - Inches(0.30), Inches(0.25),
            dataset, size=10, color=WHITE, name="Dataset",
        )
        for j, (name, desc) in enumerate(items):
            row_y = y + Inches(0.80 + j * 0.62)
            add_text(
                s, x + Inches(0.15), row_y, col_w - Inches(0.30), Inches(0.30),
                name, size=11, bold=True, color=PRIMARY, name="Code",
            )
            add_text(
                s, x + Inches(0.15), row_y + Inches(0.28), col_w - Inches(0.30), Inches(0.32),
                desc, size=9, color=MUTED,
            )
    add_footer(s, "22 governed views: profile_current_safe (default), profile_current (approved sensitive), profile_signals, profile_explain, profile_events, profile_contactability, audiences, build_performance, release_status")
    return s


def slide_build_modes(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_rect(s, 0, 0, SLIDE_W, SLIDE_H, BG)
    add_header(s, "Build Modes", "Five orchestrated modes; --build-mode is required")
    modes = [
        ("rebuild", "Full rebuild from source",
         "snapshot → ddl → maintenance → populate (5) → fill_gaps (4) → enrich (8) → personas (4) → restore → views → snapshot_core",
         "28 steps  |  ~30-40 min  |  candidate-backed publish  |  use after schema or identity-hub change"),
        ("resume_rebuild", "Pick up after a failed late-stage rebuild",
         "restore → views → snapshot_core (against existing profile_data_candidate)",
         "3 steps  |  ~4 min  |  preserves the populate/enrich work, finishes the publish"),
        ("refresh", "Daily incremental update",
         "refresh_scope → refresh → enrich (8) → personas (4) → views",
         "15 steps  |  ~15-20 min  |  3-day default lookback  |  scope-guarded (>25% hard-fails)"),
        ("reenrich", "Re-run enrich + personas only",
         "rebuild_scope → enrich (8) → personas (4) → views",
         "14 steps  |  ~10 min  |  no DDL, no populate, all bn_ids in scope"),
        ("views", "Republish consumer views only",
         "views",
         "1 step  |  ~30s  |  use after view-only SQL edit; requires populated profile_core"),
    ]
    y0 = Inches(1.55)
    row_h = Inches(1.05)
    for i, (mode, summary, steps, footer) in enumerate(modes):
        y = y0 + i * row_h
        add_rect(s, Inches(0.30), y, Inches(12.70), row_h - Inches(0.10), CARD, outline=BORDER)
        add_rect(s, Inches(0.30), y, Inches(2.40), row_h - Inches(0.10), PRIMARY)
        add_text(
            s, Inches(0.40), y + Inches(0.12), Inches(2.20), Inches(0.40),
            mode, size=16, bold=True, color=WHITE, name="Code",
        )
        add_text(
            s, Inches(0.40), y + Inches(0.55), Inches(2.20), Inches(0.30),
            summary, size=10, color=WHITE,
        )
        add_text(
            s, Inches(2.85), y + Inches(0.10), Inches(10.0), Inches(0.30),
            steps, size=10, color=BODY, name="Code",
        )
        add_text(
            s, Inches(2.85), y + Inches(0.45), Inches(10.0), Inches(0.50),
            footer, size=10, color=MUTED,
        )
    add_footer(
        s,
        "Run: python orchestrate.py --source profile_database --env prod --build-mode <mode> [--lookback 7]",
    )
    return s


def slide_quality_gates(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_rect(s, 0, 0, SLIDE_W, SLIDE_H, BG)
    add_header(
        s, "Quality Gates & Observability",
        "6 hard gates block publish; 5 soft gates warn; everything is logged",
    )
    add_text(
        s, Inches(0.40), Inches(1.55), Inches(6.20), Inches(0.40),
        "Hard gates (block publish)", size=14, bold=True, color=PRIMARY,
    )
    hard = [
        ("profile_core_unique_bn_id", "PK uniqueness on profile_core"),
        ("profile_identifiers_unique_primary", "Max 1 primary per type per bn_id"),
        ("profile_current_unique_bn_id", "View-level uniqueness check"),
        ("orphan_satellite", "Satellites must reference profile_core; populate-suppressed bn_ids excluded; soft-promoted below 0.1% per-table"),
        ("missing_parent", "engagement count within ±5% of core"),
        ("restore_coverage / refresh_safety", ">99% remap on rebuild; immutable-field check on refresh"),
    ]
    for i, (name, desc) in enumerate(hard):
        y = Inches(1.95 + i * 0.62)
        add_rect(s, Inches(0.40), y, Inches(6.20), Inches(0.55), CARD, outline=BORDER)
        add_text(
            s, Inches(0.55), y + Inches(0.04), Inches(5.95), Inches(0.30),
            name, size=10, bold=True, color=PRIMARY, name="Code",
        )
        add_text(
            s, Inches(0.55), y + Inches(0.28), Inches(5.95), Inches(0.30),
            desc, size=9, color=MUTED,
        )

    add_text(
        s, Inches(7.00), Inches(1.55), Inches(6.20), Inches(0.40),
        "Soft gates (warn only)", size=14, bold=True, color=PRIMARY,
    )
    soft = [
        ("hub_activity", "Identity hub had changes recently"),
        ("anonymous_known_count_delta", "Tier1 / tier2 distribution within 10%"),
        ("fill_rate_drift_critical / monitoring", "Per-field fill rates stable vs baseline"),
        ("profile_field_changes_populated", "Lineage rows written for tracked persona fields"),
        ("exception_spike + perf regressions", "Exception count + per-step duration baselines"),
    ]
    for i, (name, desc) in enumerate(soft):
        y = Inches(1.95 + i * 0.62)
        add_rect(s, Inches(7.00), y, Inches(6.20), Inches(0.55), CARD, outline=BORDER)
        add_text(
            s, Inches(7.15), y + Inches(0.04), Inches(5.95), Inches(0.30),
            name, size=10, bold=True, color=ACCENT, name="Code",
        )
        add_text(
            s, Inches(7.15), y + Inches(0.28), Inches(5.95), Inches(0.30),
            desc, size=9, color=MUTED,
        )

    add_text(
        s, Inches(0.40), Inches(5.85), Inches(12.5), Inches(0.40),
        "Observability tables", size=14, bold=True, color=PRIMARY,
    )
    obs = [
        ("profile_build_runs", "Run-level: status, totals, assertion summary, runtime fingerprint"),
        ("profile_build_steps", "Step-level: duration, rows, bytes processed/billed, slot-millis"),
        ("profile_publish_manifest", "Per-table promotion records from blue/green publish"),
    ]
    for i, (name, desc) in enumerate(obs):
        y = Inches(6.25 + i * 0.36)
        add_text(
            s, Inches(0.40), y, Inches(4.5), Inches(0.30),
            name, size=11, bold=True, color=PRIMARY, name="Code",
        )
        add_text(
            s, Inches(4.95), y, Inches(8.0), Inches(0.30),
            desc, size=10, color=BODY,
        )
    return s


def slide_audiences(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_rect(s, 0, 0, SLIDE_W, SLIDE_H, BG)
    add_header(
        s, "Audience Surfaces",
        "Pre-built governed views for activation, analytics, and ops",
    )
    audiences = [
        ("profile_marketing_audience", "411,916", "Marketing-safe activation",
         "communication_opt_in + email + reachability filter"),
        ("profile_audience_hcp", "305,941", "HCP targeting",
         "account_type=hcp + NPI verified"),
        ("profile_audience_patients_confirmed", "17,649", "Confirmed patients",
         "account_type=patient (zero-party + Mailchimp confirmed)"),
        ("profile_audience_caregivers", "1,241", "Caregiver targeting",
         "account_type=caregiver + relationship known"),
        ("profile_audience_high_engagement", "192,783", "High-engagement",
         "engagement_tier=high (top 5.6% of profiles)"),
        ("profile_analytics_audience", "3,460,572", "BI / analytics",
         "Aggregate-only, identifiers redacted"),
        ("profile_ops_audience", "3,460,572", "Support / CX",
         "Identifiers + key flags for individual lookup"),
        ("profile_marketing_audience (Meta)", "2,663,842", "Meta Custom Audience",
         "fbp identifier presence (whether marketable or not)"),
        ("profile_marketing_audience (Google)", "2,590,103", "Google Customer Match",
         "gcl_au identifier presence"),
    ]
    for i, (name, count, label, desc) in enumerate(audiences):
        col = i % 3
        row = i // 3
        x = Inches(0.30 + col * 4.40)
        y = Inches(1.65 + row * 1.90)
        add_rect(s, x, y, Inches(4.20), Inches(1.75), CARD, outline=BORDER)
        add_rect(s, x, y, Inches(4.20), Inches(0.40), PRIMARY)
        add_text(
            s, x + Inches(0.12), y + Inches(0.05), Inches(3.96), Inches(0.30),
            name, size=10, bold=True, color=WHITE, name="Code",
        )
        add_text(
            s, x + Inches(0.15), y + Inches(0.50), Inches(3.90), Inches(0.45),
            count, size=22, bold=True, color=ACCENT,
        )
        add_text(
            s, x + Inches(0.15), y + Inches(1.00), Inches(3.90), Inches(0.30),
            label, size=11, bold=True, color=BODY,
        )
        add_text(
            s, x + Inches(0.15), y + Inches(1.32), Inches(3.90), Inches(0.40),
            desc, size=9, color=MUTED,
        )
    return s


def slide_features_benefits(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_rect(s, 0, 0, SLIDE_W, SLIDE_H, BG)
    add_header(
        s, "Features & Benefits",
        "What downstream teams get out of the box",
    )
    items = [
        ("Single resolved profile",
         "One row per bn_id with persona, condition, engagement and signals; no manual joins required"),
        ("Cross-channel activation",
         "Email, Meta, Google audiences derived from the same identity-stitched core"),
        ("Persona explainability",
         "profile_explain surfaces source + confidence + supporting signals for every classification"),
        ("Built-in governance",
         "profile_current_safe redacts PII by default; profile_current is approved-only"),
        ("Stable bn_ids across rebuilds",
         "Persistence-aware snapshot/restore preserves app-authored fields across identity-graph re-stitches"),
        ("Observable + measurable",
         "Run + step rows + per-table publish manifest + runtime fingerprint enable cost & regression analysis"),
        ("Soft-fail performance regressions",
         "Step-level duration / bytes / slot-millis baselines flag hot-spots before users notice"),
        ("Refresh scope guards",
         "Empty scope short-circuits; >25% hard-fails; >10% warns — prevents accidental near-rebuilds"),
        ("Blue/green promotion",
         "Candidate dataset + view repointing keeps consumers consistent during table copy"),
    ]
    for i, (title, desc) in enumerate(items):
        col = i % 3
        row = i // 3
        x = Inches(0.30 + col * 4.40)
        y = Inches(1.65 + row * 1.85)
        add_rect(s, x, y, Inches(4.20), Inches(1.70), CARD, outline=BORDER)
        add_text(
            s, x + Inches(0.20), y + Inches(0.15), Inches(3.80), Inches(0.40),
            title, size=13, bold=True, color=PRIMARY,
        )
        add_rect(s, x + Inches(0.20), y + Inches(0.55), Inches(0.50), Inches(0.04), ACCENT)
        add_text(
            s, x + Inches(0.20), y + Inches(0.65), Inches(3.80), Inches(0.95),
            desc, size=11, color=BODY,
        )
    return s


# --------------------------------------------------------------------------
# Use cases + SQL slides
# --------------------------------------------------------------------------
USE_CASES = [
    {
        "num": "01",
        "title": "Person Lookup by Email",
        "summary": "Find the resolved profile for a known email",
        "context": "Use this query to resolve a known email to its full profile, including persona, condition, engagement tier, and the underlying identifier cluster.",
        "sql": """-- Resolve a known email to the full profile
SELECT
    pc.bn_id,
    pc.account_type,
    pc.preferred_condition.label AS condition,
    pc.first_name,
    pc.last_name,
    pc.country,
    pc.communication_opt_in,
    pe.engagement_tier,
    pc.profile_stage
FROM profile_data.profile_current_safe pc
LEFT JOIN profile_data.profile_engagement pe USING (bn_id)
WHERE LOWER(pc.email) = LOWER('jdoe@example.com');""",
        "tip": "Substitute profile_current for approved sensitive use; profile_current_safe redacts PII by default.",
    },
    {
        "num": "02",
        "title": "Resolve Web Visitor (bnfpvid → Profile)",
        "summary": "Turn an anonymous web cookie into a resolved profile",
        "context": "Used by personalization and on-site recommendations to look up a returning visitor's persona, condition, and engagement.",
        "sql": """-- Resolve an anonymous web cookie to a resolved profile
SELECT pc.*
FROM profile_data.profile_identifiers pi
JOIN profile_data.profile_current_safe pc USING (bn_id)
WHERE pi.identifier_type = 'bnfpvid'
  AND pi.identifier_value = 'BNFPVID-22219000-...';""",
        "tip": "Works for any identifier_type: client_id, fbp, gcl_au, npi_number, mc_euid, etc.",
    },
    {
        "num": "03",
        "title": "HCP Targeting List",
        "summary": "Verified HCPs with NPI + active engagement",
        "context": "Source of truth for HCP campaigns. Joins NPI registry data to the profile graph and filters to recent activity.",
        "sql": """-- HCP audience: verified NPI + recent activity
SELECT
    pc.bn_id,
    pc.npi_number,
    pc.specialty.label AS specialty,
    pc.credentials,
    pc.practice_state,
    pe.engagement_tier,
    pe.last_seen_web
FROM profile_data.profile_audience_hcp pc
JOIN profile_data.profile_engagement pe USING (bn_id)
WHERE pe.last_seen_web >= TIMESTAMP_SUB(CURRENT_TIMESTAMP(), INTERVAL 90 DAY)
ORDER BY pe.last_seen_web DESC;""",
        "tip": "profile_audience_hcp already filters account_type=hcp AND hcp_status=TRUE AND npi present.",
    },
    {
        "num": "04",
        "title": "Meta Custom Audience by Condition",
        "summary": "Build an fbp-keyed list for a specific condition",
        "context": "Generates the activation file for Meta Custom Audiences. Filtered to a specific preferred_condition and consent-eligible profiles.",
        "sql": """-- Meta Custom Audience: MS profiles with consent + Meta pixel
SELECT DISTINCT pi.identifier_value AS fbp
FROM profile_data.profile_marketing_audience pc
JOIN profile_data.profile_identifiers pi USING (bn_id)
WHERE pi.identifier_type = 'fbp'
  AND pc.preferred_condition.label = 'Multiple Sclerosis';""",
        "tip": "Replace 'fbp' with 'gcl_au' for Google Customer Match. profile_marketing_audience already enforces consent.",
    },
    {
        "num": "05",
        "title": "High-Engagement Newsletter Cohort",
        "summary": "Top tier readers for a newsletter targeting push",
        "context": "For a high-value campaign, target only the top-tier engaged audience, segmented by their preferred condition.",
        "sql": """-- High-engagement audience by preferred condition
SELECT
    pc.preferred_condition.label AS condition,
    COUNT(*) AS profiles,
    COUNTIF(pc.communication_opt_in) AS opted_in,
    COUNTIF(pe.email_open_count >= 5) AS active_email_readers
FROM profile_data.profile_audience_high_engagement pc
JOIN profile_data.profile_engagement pe USING (bn_id)
GROUP BY condition
ORDER BY opted_in DESC
LIMIT 20;""",
        "tip": "engagement_tier='high' is computed from VIP, member_rating, click count, sessions, and forum activity.",
    },
    {
        "num": "06",
        "title": "Persona Explainability",
        "summary": "Why is this user classified as X?",
        "context": "Debugging surface for QA and support. Shows the source, confidence, and supporting signals behind every persona / condition classification.",
        "sql": """-- One-row explainability for a profile
SELECT
    bn_id,
    account_type, account_type_source, account_type_confidence,
    preferred_condition.label AS condition,
    preferred_condition_source, preferred_condition_confidence,
    has_condition_signal_conflict,
    has_hcp_signal_conflict,
    top_content_condition,
    segment_categories
FROM profile_data.profile_explain
WHERE bn_id = 'BN_abc123';""",
        "tip": "Use has_*_conflict flags to find profiles where signals disagree with the classification.",
    },
    {
        "num": "07",
        "title": "Contactability Summary",
        "summary": "Is this user reachable, on what channel, with what consent?",
        "context": "Used by activation tools and customer support to determine reach + the main blocking reason if a profile is not reachable.",
        "sql": """-- Reachability + consent summary
SELECT
    bn_id,
    contactability_status,
    can_email_market,
    marketing_status_reason,
    has_email_marketable,
    has_paid_media_signal,
    consent_analytics, consent_advertising, consent_functional
FROM profile_data.profile_contactability
WHERE bn_id = 'BN_abc123';""",
        "tip": "marketing_status_reason explains the suppression (no_email, opted_out, no_consent, dsr_pending, etc.).",
    },
    {
        "num": "08",
        "title": "Event Timeline",
        "summary": "Unified site events + zero-party answers for one user",
        "context": "Behavioral debugging and personalization input. Merges site_events (GA4-backed) with profile_zero_party (poll/quiz answers).",
        "sql": """-- Recent events for a profile (web + zero-party unified)
SELECT
    event_timestamp,
    event_name,
    JSON_VALUE(properties, '$.source_surface') AS source,
    JSON_VALUE(properties, '$.event_label') AS label
FROM profile_data.profile_events
WHERE bn_id = 'BN_abc123'
ORDER BY event_timestamp DESC
LIMIT 100;""",
        "tip": "event_name LIKE 'zero_party_%' filters to poll/quiz answers; everything else is GA4-derived.",
    },
    {
        "num": "09",
        "title": "Coverage / Fill-Rate Monitoring",
        "summary": "Are key fields populated at the expected rate?",
        "context": "Data-quality monitoring. Track fill-rate drift on persona, demographic, and HCP fields against the latest baseline.",
        "sql": """-- Fill-rate snapshot
SELECT cohort, field_name, populated_n, total_n, fill_rate_pct
FROM profile_data.profile_coverage
WHERE cohort IN ('all', 'patient', 'hcp')
ORDER BY cohort, fill_rate_pct;""",
        "tip": "Compare against profile_ops.profile_build_runs.metadata.fill_rate_baseline for drift detection.",
    },
    {
        "num": "10",
        "title": "Build Performance & Cost",
        "summary": "Which steps are slow, expensive, or regressing?",
        "context": "Operations and tuning. Surfaces per-step duration, bytes processed/billed, slot-millis, and per-step regression vs same-mode baseline.",
        "sql": """-- Latest build per mode with per-step performance
SELECT
    build_id, mode, step_name,
    duration_seconds AS secs,
    ROUND(step_total_bytes_processed / POW(10, 9), 2) AS gb_scanned,
    ROUND(step_total_slot_millis / 1000.0, 1) AS slot_seconds
FROM profile_data.profile_build_performance
WHERE is_latest_for_mode
ORDER BY build_started_at DESC, secs DESC;""",
        "tip": "step_performance_regression soft-fails when a step exceeds recent same-mode baseline; check assertion_summary.",
    },
    {
        "num": "11",
        "title": "Release & Publish Health",
        "summary": "Current release state + per-table promotion progress",
        "context": "Operations dashboard. Shows the current published release, the source dataset behind production views, and any in-flight publish.",
        "sql": """-- Release health (one row per build)
SELECT * FROM profile_data.profile_release_status;

-- Per-table promotion records for the latest publish
SELECT table_name, status, started_at, completed_at,
       row_count_target, row_count_source
FROM profile_ops.profile_publish_manifest
WHERE build_id = (SELECT MAX(build_id) FROM profile_ops.profile_publish_manifest)
ORDER BY started_at;""",
        "tip": "Use profile_publish_manifest.status for resume_publish recovery: rows in 'copying' or 'failed' are the targets.",
    },
    {
        "num": "12",
        "title": "Data Quality Exceptions",
        "summary": "Conflicts and anomalies surfaced by gates",
        "context": "Data-quality triage. profile_exceptions surfaces conflicts like HCP-without-NPI, patient-without-condition, high-tier-no-sessions.",
        "sql": """-- Top exception types
SELECT exception_type,
       COUNT(*) AS n_rows,
       COUNT(DISTINCT bn_id) AS n_profiles
FROM profile_data.profile_exceptions
GROUP BY exception_type
ORDER BY n_rows DESC;""",
        "tip": "Use exception_type to filter to a specific issue (hcp_no_npi, subtype_without_stage, etc.).",
    },
]


def slide_use_cases(prs, sql_slide_indices):
    """Use Cases navigator slide — 12 cards, each hyperlinks to its SQL slide.

    sql_slide_indices: list of pptx slide objects to link to (one per use case).
    """
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_rect(s, 0, 0, SLIDE_W, SLIDE_H, BG)
    add_header(
        s, "Use Cases",
        "Click any example to view the SQL  —  all queries run against profile_data in BigQuery",
    )
    for i, uc in enumerate(USE_CASES):
        col = i // 6
        row = i % 6
        x = Inches(0.28 + col * 6.66)
        y = Inches(1.55 + row * 0.90)
        card = add_rect(s, x, y, Inches(6.10), Inches(0.84), CARD, outline=BORDER)
        # Wire click action on the card itself
        try:
            card.click_action.target_slide = sql_slide_indices[i]
        except Exception:
            pass
        add_text(
            s, x + Inches(0.12), y + Inches(0.10),
            Inches(0.50), Inches(0.56),
            uc["num"], size=18, bold=True, color=ACCENT, anchor=MSO_ANCHOR.MIDDLE,
        )
        add_text(
            s, x + Inches(0.60), y + Inches(0.10),
            Inches(5.10), Inches(0.32),
            uc["title"], size=12, bold=True, color=PRIMARY,
        )
        add_text(
            s, x + Inches(0.60), y + Inches(0.46),
            Inches(5.10), Inches(0.30),
            uc["summary"], size=9, color=MUTED,
        )
        arrow = add_text(
            s, x + Inches(5.78), y + Inches(0.28),
            Inches(0.24), Inches(0.28),
            "▶", size=14, bold=True, color=ACCENT,
        )
        try:
            arrow.click_action.target_slide = sql_slide_indices[i]
        except Exception:
            pass
    return s


def slide_sql(prs, uc, back_target):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_rect(s, 0, 0, SLIDE_W, SLIDE_H, BG)
    add_header(s, f"SQL: {uc['title']}", uc["summary"])

    add_text(
        s, Inches(0.40), Inches(1.52), Inches(11.00), Inches(0.48),
        uc["context"], size=11, color=BODY,
    )
    # Back button
    back = add_rect(s, Inches(11.65), Inches(1.48), Inches(1.48), Inches(0.38), ACCENT)
    try:
        back.click_action.target_slide = back_target
    except Exception:
        pass
    back_txt = add_text(
        s, Inches(11.65), Inches(1.48), Inches(1.48), Inches(0.38),
        "◀  All Examples", size=10, bold=True, color=WHITE,
        align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
    )
    try:
        back_txt.click_action.target_slide = back_target
    except Exception:
        pass

    # Code block
    code_x = Inches(0.40)
    code_y = Inches(2.06)
    code_w = Inches(12.50)
    code_h = Inches(4.50)
    add_rect(s, code_x, code_y, code_w, code_h, CODE_BG)
    code_tb = s.shapes.add_textbox(
        code_x + Inches(0.15), code_y + Inches(0.10),
        code_w - Inches(0.30), code_h - Inches(0.20),
    )
    code_tb.name = "Code"
    tf = code_tb.text_frame
    tf.margin_left = Pt(0)
    tf.margin_right = Pt(0)
    tf.margin_top = Pt(0)
    tf.margin_bottom = Pt(0)
    tf.word_wrap = True
    for i, line in enumerate(uc["sql"].split("\n")):
        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        para.alignment = PP_ALIGN.LEFT
        run = para.add_run()
        run.text = line if line else " "
        run.font.size = Pt(13)
        run.font.name = "Consolas"
        # Comment lines get muted color, SQL keywords get accent
        stripped = line.lstrip()
        if stripped.startswith("--"):
            run.font.color.rgb = RGBColor(0x90, 0xA4, 0xAE)
            run.font.italic = True
        else:
            run.font.color.rgb = CODE_TXT

    # Tip box at the bottom
    add_rect(s, Inches(0.40), Inches(6.70), Inches(12.50), Inches(0.55), CARD, outline=BORDER)
    add_text(
        s, Inches(0.55), Inches(6.78), Inches(0.5), Inches(0.40),
        "TIP", size=10, bold=True, color=ACCENT, anchor=MSO_ANCHOR.MIDDLE,
    )
    add_text(
        s, Inches(1.10), Inches(6.78), Inches(11.7), Inches(0.40),
        uc["tip"], size=10, color=BODY, anchor=MSO_ANCHOR.MIDDLE,
    )
    return s


def slide_operations(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_rect(s, 0, 0, SLIDE_W, SLIDE_H, BG)
    add_header(s, "Operations", "Daily refresh, full rebuild, recovery from partial state")
    blocks = [
        ("Daily refresh (most common)",
         "python orchestrate.py --source profile_database --env prod --build-mode refresh\n"
         "# scope guard auto-fails if scope > 25%; warn if > 10%\n"
         "# empty scope short-circuits as completed_no_changes",
         "Run via Cloud Scheduler / Airflow"),
        ("Full rebuild",
         "python orchestrate.py --source profile_database --env prod --build-mode rebuild\n"
         "# blue/green publish through profile_data_candidate\n"
         "# 28 steps; ~30-40 min wall-clock; 6 hard + 5 soft assertions",
         "After identity-hub regeneration or schema change"),
        ("Resume after partial publish",
         "python orchestrate.py --source profile_database --env prod --build-mode resume_rebuild\n"
         "# picks up restore → views → snapshot_core against existing\n"
         "# profile_data_candidate; preserves the 30 min populate work",
         "When rebuild fails at restore / views / snapshot_core"),
        ("Re-enrich only",
         "python orchestrate.py --source profile_database --env prod --build-mode reenrich\n"
         "# rebuild_scope sentinel + enrich + personas + views\n"
         "# 14 steps; no DDL, no populate",
         "After enrich/persona SQL change"),
        ("Views-only republish",
         "python orchestrate.py --source profile_database --env prod --build-mode views\n"
         "# rebuilds the 22 governed views in place\n"
         "# requires profile_core to exist (errors with clear message if not)",
         "After view-only SQL edit"),
    ]
    for i, (title, code, footer) in enumerate(blocks):
        y = Inches(1.55 + i * 1.15)
        add_rect(s, Inches(0.30), y, Inches(12.70), Inches(1.05), CARD, outline=BORDER)
        add_text(
            s, Inches(0.50), y + Inches(0.05), Inches(12.5), Inches(0.30),
            title, size=12, bold=True, color=PRIMARY,
        )
        # Code block
        for j, line in enumerate(code.split("\n")):
            add_text(
                s, Inches(0.50), y + Inches(0.34 + j * 0.20),
                Inches(12.5), Inches(0.20),
                line, size=10,
                color=BODY if not line.lstrip().startswith("#") else MUTED,
                name="Code",
            )
    return s


def slide_roadmap(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_rect(s, 0, 0, SLIDE_W, SLIDE_H, BG)
    add_header(s, "Roadmap", "What's next for the profile database")
    items = [
        ("v6.5 (in progress)",
         "Correctness hardening",
         "• Refresh-scope contract + audit gate\n"
         "• Lineage coverage matrix + profile_explain guardrails\n"
         "• Operator failure playbook + check-gated generation\n"
         "• Release manifest + resume_publish + post-copy verification\n"
         "• Unit tests for _rewrite_internal_datasets"),
        ("v6.6",
         "Refresh isolation + parity",
         "• Candidate-backed refresh for large scopes\n"
         "• Stratified sampled shadow rebuild/refresh diff\n"
         "• Stronger automation around final view repoint / publish recovery\n"
         "• Promote scope_predicate_audit.py to hard gate"),
        ("v6.7",
         "App-owned source + retire snapshot/restore",
         "• enrich_app_attributes reading sso_app_data.user_attributes directly\n"
         "• One rebuild with both snapshot/restore and app-source active; diff\n"
         "• Retire snapshot/restore after one clean cycle\n"
         "• Dependency: bionews_uk SSO coverage (today: ~300 profiles)"),
        ("Deferred",
         "Awaiting cross-team alignment",
         "• Canonical consent fact (legal + marketing sign-off needed)\n"
         "• Engagement table consolidation\n"
         "• Per-channel marketing eligibility gates"),
    ]
    for i, (label, title, body) in enumerate(items):
        col = i % 2
        row = i // 2
        x = Inches(0.40 + col * 6.40)
        y = Inches(1.65 + row * 2.85)
        add_rect(s, x, y, Inches(6.20), Inches(2.65), CARD, outline=BORDER)
        add_rect(s, x, y, Inches(6.20), Inches(0.55), PRIMARY)
        add_text(
            s, x + Inches(0.20), y + Inches(0.05), Inches(6.0), Inches(0.30),
            label, size=12, bold=True, color=ACCENT,
        )
        add_text(
            s, x + Inches(0.20), y + Inches(0.30), Inches(6.0), Inches(0.30),
            title, size=14, bold=True, color=WHITE,
        )
        add_text(
            s, x + Inches(0.20), y + Inches(0.70), Inches(6.0), Inches(1.90),
            body, size=11, color=BODY,
        )
    return s


def slide_thank_you(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_rect(s, 0, 0, SLIDE_W, SLIDE_H, PRIMARY)
    add_rect(s, 0, Inches(5.6), SLIDE_W, Inches(1.9), DEEP)
    add_rect(s, 0, Inches(5.6), SLIDE_W, Inches(0.06), ACCENT)
    add_text(
        s, Inches(0.6), Inches(2.2), Inches(12.0), Inches(0.6),
        "BioNews Profile Database", size=28, bold=True, color=WHITE,
    )
    add_text(
        s, Inches(0.6), Inches(3.0), Inches(12.0), Inches(1.2),
        "3.46M profiles. One identity. Every channel.",
        size=44, bold=True, color=WHITE,
    )
    add_text(
        s, Inches(0.6), Inches(5.85), Inches(12.0), Inches(0.5),
        "Default surface: profile_data.profile_current_safe",
        size=18, color=WHITE,
    )
    add_text(
        s, Inches(0.6), Inches(6.4), Inches(12.0), Inches(0.4),
        "Operator runbook: docs/PROFILE_DATABASE_OPERATOR_GUIDE.md",
        size=14, color=MUTED,
    )
    return s


# --------------------------------------------------------------------------
# Driver
# --------------------------------------------------------------------------
def build():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    slide_cover(prs)
    slide_agenda(prs)
    slide_what_is(prs)
    slide_scale(prs)
    slide_persona_mix(prs)
    slide_architecture(prs)
    slide_output_tables(prs)
    slide_build_modes(prs)
    slide_quality_gates(prs)
    slide_audiences(prs)
    slide_features_benefits(prs)
    # First create empty SQL slides so we have target references for hyperlinks
    sql_slides = [None] * len(USE_CASES)
    # We need the use cases nav slide and SQL slides; add nav first as placeholder, then
    # build SQL slides, then go back and wire the nav slide. python-pptx click_action
    # requires the target slide to exist BEFORE the source shape is linked. So:
    # 1. Build SQL slides first (with placeholder back-link shapes).
    # 2. Build nav slide last, hyperlink to SQL slides.
    # Then patch the SQL slides' back-link to point to the nav slide.
    # Simplest: build SQL slides with a deferred back-target, then nav, then patch.
    # Actually simpler: nav slide must be created first (it's earlier in deck), but
    # add_slide returns the slide object and click_action can target a slide that
    # doesn't yet exist if we set it later. Best approach: create nav slide as placeholder,
    # create SQL slides referencing it, then go back and add hyperlinks on nav cards.

    # Create the nav slide first as a stub (cards without click actions) so SQL slides
    # can target it for back-links. Then patch the nav slide's cards to point at SQL slides.
    nav_slide_idx_before = len(prs.slides)

    # Create nav slide (with no working hyperlinks yet)
    nav_slide = prs.slides.add_slide(prs.slide_layouts[6])
    # Defer drawing the cards on nav slide; we'll fill it in after SQL slides exist.

    # Now create SQL slides, each with a back-link to nav_slide
    for i, uc in enumerate(USE_CASES):
        sql_slides[i] = slide_sql(prs, uc, back_target=nav_slide)

    # Now go back and draw the nav slide's card grid with hyperlinks to sql_slides
    add_rect(nav_slide, 0, 0, SLIDE_W, SLIDE_H, BG)
    add_header(
        nav_slide, "Use Cases",
        "Click any example to view the SQL  —  all queries run against profile_data in BigQuery",
    )
    for i, uc in enumerate(USE_CASES):
        col = i // 6
        row = i % 6
        x = Inches(0.28 + col * 6.66)
        y = Inches(1.55 + row * 0.90)
        card = add_rect(nav_slide, x, y, Inches(6.10), Inches(0.84), CARD, outline=BORDER)
        try:
            card.click_action.target_slide = sql_slides[i]
        except Exception:
            pass
        add_text(
            nav_slide, x + Inches(0.12), y + Inches(0.10),
            Inches(0.50), Inches(0.56),
            uc["num"], size=18, bold=True, color=ACCENT, anchor=MSO_ANCHOR.MIDDLE,
        )
        add_text(
            nav_slide, x + Inches(0.60), y + Inches(0.10),
            Inches(5.10), Inches(0.32),
            uc["title"], size=12, bold=True, color=PRIMARY,
        )
        add_text(
            nav_slide, x + Inches(0.60), y + Inches(0.46),
            Inches(5.10), Inches(0.30),
            uc["summary"], size=9, color=MUTED,
        )
        arrow = add_text(
            nav_slide, x + Inches(5.78), y + Inches(0.28),
            Inches(0.24), Inches(0.28),
            "▶", size=14, bold=True, color=ACCENT,
        )
        try:
            arrow.click_action.target_slide = sql_slides[i]
        except Exception:
            pass

    slide_operations(prs)
    slide_thank_you(prs)

    out = Path(__file__).resolve().parent / "Profile_Database_Presentation.pptx"
    prs.save(out)
    print(f"[OK] Wrote {out}  ({len(prs.slides)} slides)")
    return out


if __name__ == "__main__":
    build()
