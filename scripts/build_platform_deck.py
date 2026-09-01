#!/usr/bin/env python3
"""
Build exports/bionews_data_platform_bi.pptx -- the BI-facing deck for the
identity hub and profile database as one platform.

The two systems are usually described separately, which hides the thing that
matters: the hub works out WHO someone is, the profile database records WHAT we
know about them, and bn_id is the join between them. Neither is usable alone.
This deck tells that story in order, then hands over the eight reference
documents at the end.

Audience: people who will write queries. Depth is spent where it prevents a
wrong number, not on architecture for its own sake.

Every figure is production-verified as of AS_OF below and matches the .docx
set. Document links are relative, so the deck must ship in the same folder.

  python scripts/build_platform_deck.py
"""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

REPO = Path(__file__).resolve().parents[1]
EXPORTS = REPO / "exports"
OUT = EXPORTS / "bionews_data_platform_bi.pptx"
AS_OF = "24 August 2026"

NAVY = RGBColor(0x0D, 0x47, 0xA1)
DEEP = RGBColor(0x0A, 0x2A, 0x5E)
TEAL = RGBColor(0x1A, 0x73, 0xE8)
SKY = RGBColor(0xBF, 0xD4, 0xF2)
GOLD = RGBColor(0xFB, 0xBC, 0x04)
AMBER_BG = RGBColor(0xFF, 0xF8, 0xE1)
AMBER_TX = RGBColor(0x7A, 0x5A, 0x00)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LGRAY = RGBColor(0xF8, 0xF9, 0xFA)
BORDER = RGBColor(0xDA, 0xDC, 0xE0)
DGRAY = RGBColor(0x21, 0x21, 0x21)
MGRAY = RGBColor(0x5F, 0x63, 0x68)
GREEN = RGBColor(0x34, 0xA8, 0x53)
RED = RGBColor(0xC0, 0x00, 0x00)
REDBG = RGBColor(0xFF, 0xF3, 0xF3)
PURPLE = RGBColor(0x7B, 0x1F, 0xA2)
SLATE = RGBColor(0x90, 0xA4, 0xAE)

# The eight documents, in the order they should be handed over.
DOCS = [
    (
        "identity_hub_overview.docx",
        "Identity Hub -- Overview",
        NAVY,
        "Why the graph exists and how it decides two records are the same person.",
    ),
    (
        "identity_hub_bi_queries.docx",
        "Identity Hub -- BI Queries",
        NAVY,
        "Fourteen runnable queries, each with the trap it avoids stated first.",
    ),
    (
        "identity_hub_table_reference.docx",
        "Identity Hub -- Table Reference",
        NAVY,
        "The eight hub tables, with confidence and cluster-health scoring in full.",
    ),
    (
        "identity_hub_data_dictionary.docx",
        "Identity Hub -- Data Dictionary",
        NAVY,
        "Every hub column, generated from the DDL and the pipeline config.",
    ),
    (
        "profile_db_overview.docx",
        "Profile Database -- Overview",
        PURPLE,
        "What gets recorded once the graph has worked out who someone is.",
    ),
    (
        "profile_db_bi_queries.docx",
        "Profile Database -- BI Queries",
        PURPLE,
        "Fifteen queries plus the rules that stop a number being quoted wrongly.",
    ),
    (
        "profile_db_table_reference.docx",
        "Profile Database -- Table Reference",
        PURPLE,
        "Tables grouped by theme, leading with the views analysts actually query.",
    ),
    (
        "profile_db_data_dictionary.docx",
        "Profile Database -- Data Dictionary",
        PURPLE,
        "Every table AND view column -- the only document covering the views.",
    ),
]


# ── primitives ────────────────────────────────────────────────────────────────


def blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def bg(slide, rgb):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = rgb


def rect(slide, l, t, w, h, fill, line=None, line_pt=0.75, shape=MSO_SHAPE.RECTANGLE):
    sh = slide.shapes.add_shape(shape, Inches(l), Inches(t), Inches(w), Inches(h))
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill
    if line is not None:
        sh.line.color.rgb = line
        sh.line.width = Pt(line_pt)
    else:
        sh.line.fill.background()
    sh.shadow.inherit = False
    return sh


def chip(slide, s, l, t, w, h, fill, color=WHITE, size=9.5, bold=True):
    sh = rect(slide, l, t, w, h, fill, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    tf = sh.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = tf.margin_right = Inches(0.04)
    tf.margin_top = tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = s
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = color
    return sh


def arrow(slide, l, t, w, h, fill=SLATE):
    return rect(slide, l, t, w, h, fill, shape=MSO_SHAPE.RIGHT_ARROW)


def text(
    slide,
    s,
    l,
    t,
    w,
    h,
    size=14,
    bold=False,
    color=DGRAY,
    align=PP_ALIGN.LEFT,
    italic=False,
    font=None,
    link=None,
    anchor=None,
):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    if anchor:
        tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = s
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.italic = italic
    r.font.color.rgb = color
    if font:
        r.font.name = font
    if link:
        # Relative address: resolves wherever the folder is copied. An absolute
        # path would work only on the machine that built the deck.
        r.hyperlink.address = link
        r.font.underline = True
    return tb


def bullets(slide, items, l, t, w, h, size=12.5, gap=8, bullet="—"):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(gap)
        if isinstance(item, tuple):
            lead, body = item
            r = p.add_run()
            r.text = f"{bullet} {lead}  "
            r.font.size = Pt(size)
            r.font.bold = True
            r.font.color.rgb = DGRAY
            r2 = p.add_run()
            r2.text = body
            r2.font.size = Pt(size)
            r2.font.color.rgb = MGRAY
        else:
            r = p.add_run()
            r.text = f"{bullet} {item}"
            r.font.size = Pt(size)
            r.font.color.rgb = DGRAY
    return tb


def header(slide, kicker, title, subtitle=None, accent=TEAL):
    rect(slide, 0, 0, 13.33, 1.22, NAVY)
    if kicker:
        text(
            slide,
            kicker.upper(),
            0.45,
            0.11,
            8,
            0.24,
            size=9.5,
            bold=True,
            color=accent,
        )
    text(slide, title, 0.45, 0.32, 12.4, 0.52, size=25, bold=True, color=WHITE)
    if subtitle:
        text(slide, subtitle, 0.45, 0.85, 12.4, 0.32, size=12, color=SKY)
    rect(slide, 0, 1.22, 13.33, 0.05, accent)


def footer(slide, n, note=None):
    text(
        slide,
        note or f"BioNews data platform  |  {AS_OF}",
        0.45,
        7.05,
        9,
        0.28,
        size=9,
        color=MGRAY,
    )
    text(
        slide,
        str(n),
        12.55,
        7.05,
        0.35,
        0.28,
        size=9,
        color=MGRAY,
        align=PP_ALIGN.RIGHT,
    )


def takeaway(slide, s, colour=TEAL, y=6.35, band=None):
    rect(slide, 0.5, y, 12.35, 0.6, band or RGBColor(0xE8, 0xF0, 0xFE), line=colour)
    text(slide, s, 0.75, y + 0.13, 11.9, 0.4, size=12.5, bold=True, color=DGRAY)


def seealso(slide, label, fname, l, t, w=5.6):
    """A "more detail here" line linking to one of the eight documents."""
    text(slide, label, l, t, 2.3, 0.28, size=10.5, color=MGRAY)
    text(slide, fname, l + 1.5, t, w, 0.28, size=10.5, bold=True, color=TEAL,
         font="Consolas", link=fname)


# ── slides ────────────────────────────────────────────────────────────────────


def s_title(prs):
    s = blank(prs)
    bg(s, NAVY)
    rect(s, 0, 0, 13.33, 7.5, NAVY)
    rect(s, 0, 4.75, 13.33, 2.75, DEEP)
    rect(s, 0.65, 3.3, 2.4, 0.07, TEAL)
    text(
        s,
        "The BioNews Data Platform",
        0.65,
        1.15,
        12,
        0.95,
        size=42,
        bold=True,
        color=WHITE,
    )
    text(
        s,
        "Identity Hub and Profile Database -- one system, two halves",
        0.65,
        2.35,
        12,
        0.5,
        size=20,
        color=TEAL,
    )
    text(
        s,
        "The hub works out WHO someone is. The profile database records WHAT "
        "we know about them. bn_id is the join.",
        0.65,
        3.55,
        11.6,
        0.6,
        size=14,
        color=SKY,
    )
    text(s, "For the BI team", 0.65, 5.15, 6, 0.35, size=14, bold=True, color=WHITE)
    text(
        s,
        f"All figures production-verified, {AS_OF}",
        0.65,
        5.55,
        8,
        0.32,
        size=11,
        color=MGRAY,
    )
    text(
        s,
        "Underlined document names throughout are clickable. Links are relative -- "
        "keep this deck alongside the eight .docx files.",
        0.65,
        6.15,
        11.6,
        0.5,
        size=10.5,
        italic=True,
        color=SKY,
    )
    text(
        s, "CONFIDENTIAL -- Internal Use Only", 0.65, 7.0, 8, 0.3, size=9.5, color=MGRAY
    )
    return s


def s_two_halves(prs, n):
    s = blank(prs)
    bg(s, LGRAY)
    header(
        s,
        "the shape of it",
        "One platform, two halves",
        "Described separately they look like two projects. They are one pipeline with a join in the middle.",
    )
    # hub panel
    rect(s, 0.5, 1.55, 5.6, 3.5, WHITE, line=BORDER)
    rect(s, 0.5, 1.55, 5.6, 0.07, NAVY)
    text(s, "IDENTITY HUB", 0.75, 1.72, 4.5, 0.3, size=12, bold=True, color=NAVY)
    text(s, "Answers WHO", 0.75, 2.02, 4.5, 0.35, size=19, bold=True, color=DGRAY)
    bullets(
        s,
        [
            "Takes every identifier we ever see",
            "Decides which belong to the same person",
            "Issues one durable bn_id per person",
            "Keeps the evidence for every decision",
        ],
        0.75,
        2.5,
        5.1,
        1.8,
        size=12,
    )
    text(
        s,
        "34.2M identifier rows  |  7.51M people",
        0.75,
        4.55,
        5.1,
        0.3,
        size=11,
        bold=True,
        color=NAVY,
    )
    # profile panel
    rect(s, 7.23, 1.55, 5.6, 3.5, WHITE, line=BORDER)
    rect(s, 7.23, 1.55, 5.6, 0.07, PURPLE)
    text(s, "PROFILE DATABASE", 7.48, 1.72, 4.5, 0.3, size=12, bold=True, color=PURPLE)
    text(
        s,
        "Answers WHAT ABOUT THEM",
        7.48,
        2.02,
        5.1,
        0.35,
        size=19,
        bold=True,
        color=DGRAY,
    )
    bullets(
        s,
        [
            "One row per person, keyed on bn_id",
            "Condition, role, consent, engagement",
            "Merges eight source systems",
            "Rebuilt daily, gated before publish",
        ],
        7.48,
        2.5,
        5.1,
        1.8,
        size=12,
    )
    text(
        s,
        "7.50M profiles  |  131 columns on the core",
        7.48,
        4.55,
        5.1,
        0.3,
        size=11,
        bold=True,
        color=PURPLE,
    )
    # join
    chip(s, "bn_id", 6.15, 3.05, 1.0, 0.42, TEAL, size=12)
    text(
        s, "the join", 6.15, 3.5, 1.0, 0.25, size=9, color=MGRAY, align=PP_ALIGN.CENTER
    )
    takeaway(
        s,
        "You cannot read the profile database without knowing what a bn_id is. "
        "That is why the hub comes first.",
    )
    seealso(s, "Hub concepts:", "identity_hub_overview.docx", 0.5, 5.2)
    seealso(s, "Profile concepts:", "profile_db_overview.docx", 6.9, 5.2, 4.7)
    footer(s, n)
    return s


def s_problem(prs, n):
    s = blank(prs)
    bg(s, LGRAY)
    header(
        s,
        "the problem",
        "One person arrives as many records",
        "Nobody introduces themselves the same way twice. Every system stores a different fragment.",
    )
    frag = [
        ("email", 1.05),
        ("browser cookie", 1.62),
        ("newsletter ID", 2.19),
        ("forum login", 2.76),
        ("survey response", 3.33),
        ("ad click ID", 3.90),
    ]
    text(
        s,
        "WHAT WE ACTUALLY RECEIVE",
        0.6,
        1.55,
        4,
        0.26,
        size=10,
        bold=True,
        color=MGRAY,
    )
    for label, y in frag:
        chip(s, label, 0.6, y, 2.5, 0.44, WHITE, color=DGRAY, size=10.5, bold=False)
        rect(s, 0.6, y, 0.06, 0.44, SLATE)
    arrow(s, 3.45, 2.35, 1.5, 0.75, TEAL)
    text(
        s,
        "identity\nresolution",
        3.45,
        3.18,
        1.5,
        0.5,
        size=10,
        bold=True,
        color=TEAL,
        align=PP_ALIGN.CENTER,
    )
    rect(s, 5.35, 1.55, 3.5, 2.85, WHITE, line=TEAL, line_pt=1.5)
    rect(s, 5.35, 1.55, 3.5, 0.07, TEAL)
    text(s, "ONE PERSON", 5.6, 1.75, 3.0, 0.3, size=11, bold=True, color=TEAL)
    text(
        s,
        "bn_id",
        5.6,
        2.05,
        3.0,
        0.4,
        size=22,
        bold=True,
        color=DGRAY,
        font="Consolas",
    )
    bullets(
        s,
        [
            "6 identifiers",
            "3 conditions of interest",
            "1 consent state",
            "1 engagement history",
        ],
        5.6,
        2.6,
        3.0,
        1.6,
        size=11,
    )
    text(s, "WHY IT MATTERS", 9.35, 1.55, 3.5, 0.26, size=10, bold=True, color=MGRAY)
    rect(s, 9.35, 1.85, 3.5, 2.55, WHITE, line=BORDER)
    bullets(
        s,
        [
            "Counting records counts browsers, not people",
            "The same person is emailed twice",
            "Engagement looks lower than it is",
            "Consent applies to a fragment, not a person",
        ],
        9.55,
        2.0,
        3.1,
        2.3,
        size=11,
    )
    takeaway(
        s,
        "Every audience number in this deck depends on this step being right. "
        "Get it wrong and everything downstream is wrong in the same direction.",
    )
    footer(s, n)
    return s


def s_identifiers(prs, n):
    s = blank(prs)
    bg(s, LGRAY)
    header(
        s,
        "inputs",
        "Where identifiers come from",
        "22 match rules across web, email, forms, advertising and CRM. Production counts, distinct people.",
    )
    rows = [
        ("bnfpvid", 6_839_496, "first-party browser ID", NAVY),
        ("client_id", 6_389_175, "GA4 client", NAVY),
        ("gcl_au", 5_199_655, "Google Ads click", SLATE),
        ("dmd_tag", 4_243_536, "advertising tag", SLATE),
        ("aim_tag_id", 3_960_515, "AIM clickstream", SLATE),
        ("fbp", 2_367_916, "Meta pixel", SLATE),
        ("email", 790_188, "the strongest person anchor", GREEN),
        ("mc_euid", 441_423, "Mailchimp contact", TEAL),
    ]
    top = max(r[1] for r in rows)
    y = 1.5
    for name, cnt, desc, colour in rows:
        w = 6.6 * (cnt / top)
        rect(s, 3.15, y, w, 0.44, colour)
        text(
            s,
            name,
            0.6,
            y + 0.06,
            2.4,
            0.32,
            size=12,
            bold=True,
            color=DGRAY,
            font="Consolas",
        )
        text(s, f"{cnt:,}", 3.3, y + 0.08, 2.2, 0.3, size=11, bold=True, color=WHITE)
        text(s, desc, 10.0, y + 0.08, 2.9, 0.3, size=10.5, color=MGRAY)
        y += 0.6
    takeaway(
        s,
        "Browser identifiers outnumber emails nine to one. That gap IS the "
        "known-versus-unknown problem -- most of what we see is a browser, not a person.",
        colour=GOLD,
        band=AMBER_BG,
    )
    footer(s, n)
    return s


def s_matching(prs, n):
    s = blank(prs)
    bg(s, LGRAY)
    header(
        s,
        "how it works",
        "How the hub decides two records match",
        "Every link carries a confidence score. Only links at or above 0.80 are allowed to merge people.",
    )
    text(s, "CONFIDENCE SCALE", 0.6, 1.5, 5, 0.26, size=10, bold=True, color=MGRAY)
    rect(s, 0.6, 1.85, 12.2, 0.55, RGBColor(0xEC, 0xEF, 0xF1))
    rect(s, 0.6, 1.85, 12.2 * 0.80, 0.55, RGBColor(0xFF, 0xCD, 0xD2))
    rect(s, 0.6 + 12.2 * 0.80, 1.85, 12.2 * 0.20, 0.55, RGBColor(0xC8, 0xE6, 0xC9))
    rect(s, 0.6 + 12.2 * 0.80 - 0.02, 1.72, 0.05, 0.82, RED)
    text(s, "0.00", 0.6, 2.45, 1, 0.25, size=9, color=MGRAY)
    text(
        s,
        "0.80  stitch threshold",
        0.6 + 12.2 * 0.80 - 1.1,
        2.45,
        2.6,
        0.25,
        size=9.5,
        bold=True,
        color=RED,
    )
    text(
        s,
        "stored as evidence, never merges people",
        1.2,
        1.99,
        5,
        0.3,
        size=11,
        bold=True,
        color=RGBColor(0xB7, 0x1C, 0x1C),
    )
    text(
        s,
        "merges",
        10.9,
        1.99,
        1.8,
        0.3,
        size=11,
        bold=True,
        color=RGBColor(0x1B, 0x5E, 0x20),
    )
    cards = [
        (
            "DETERMINISTIC",
            "Same email on two records. Same SSO key. Same form row.",
            "1.00",
            GREEN,
        ),
        (
            "PROBABILISTIC",
            "Same device, same IP, close in time. Scored, capped per rule.",
            "0.50 - 0.85",
            TEAL,
        ),
        (
            "REJECTED",
            "Shared workstation, bot traffic, an identifier seen on 125+ people.",
            "excluded",
            SLATE,
        ),
    ]
    x = 0.6
    for title, body, score, colour in cards:
        rect(s, x, 3.1, 3.95, 1.85, WHITE, line=BORDER)
        rect(s, x, 3.1, 3.95, 0.06, colour)
        text(s, title, x + 0.22, 3.28, 3.5, 0.28, size=11, bold=True, color=colour)
        text(s, score, x + 0.22, 3.56, 3.5, 0.35, size=16, bold=True, color=DGRAY)
        text(s, body, x + 0.22, 3.98, 3.55, 0.85, size=11, color=MGRAY)
        x += 4.13
    text(
        s,
        "Evidence is kept either way. bn_id_hub holds all 75,405,402 links, "
        "including the ones too weak to act on -- which is why counting edges is "
        "not counting people.",
        0.6,
        5.15,
        12.2,
        0.5,
        size=11.5,
        italic=True,
        color=MGRAY,
    )
    takeaway(
        s,
        "A bn_id is a cluster of identifiers, not a user ID. It can absorb "
        "another cluster, and it can be superseded -- always resolve through bn_id_xref.",
    )
    seealso(s, "Scoring in full:", "identity_hub_table_reference.docx", 0.6, 5.72)
    footer(s, n)
    return s


def s_hub_output(prs, n):
    s = blank(prs)
    bg(s, LGRAY)
    header(
        s,
        "output",
        "What the hub gives you",
        "One table does almost all the BI work. The rest is evidence and machinery.",
    )
    main = [
        (
            "bn_id_xref",
            "one row per identifier, mapped to its person",
            "The table you join to. Filter is_bot = FALSE in any people count.",
            TEAL,
        ),
        (
            "bn_id_persistence",
            "where a superseded bn_id now points",
            "Needed whenever you hold a bn_id from an earlier extract.",
            NAVY,
        ),
    ]
    y = 1.5
    for name, sub, why, colour in main:
        rect(s, 0.6, y, 12.2, 1.15, WHITE, line=BORDER)
        rect(s, 0.6, y, 0.07, 1.15, colour)
        text(
            s,
            name,
            0.9,
            y + 0.14,
            3.4,
            0.32,
            size=15,
            bold=True,
            color=colour,
            font="Consolas",
        )
        text(s, sub, 0.9, y + 0.5, 4.6, 0.3, size=11, color=MGRAY)
        text(s, why, 5.9, y + 0.3, 6.6, 0.6, size=12, color=DGRAY)
        y += 1.3
    text(
        s,
        "ALSO PRESENT, RARELY NEEDED FOR BI",
        0.6,
        4.2,
        6,
        0.26,
        size=10,
        bold=True,
        color=MGRAY,
    )
    others = [
        ("bn_id_hub", "every link, incl. rejected"),
        ("bn_id_identity_changes", "merge audit trail"),
        ("bn_id_neighbors", "cluster shape"),
        ("bn_id_node_index", "rebuild machinery"),
        ("bn_id_metrics", "run stats"),
        ("bn_id_manifest", "which run is live"),
    ]
    x = 0.6
    for name, sub in others:
        chip(
            s,
            name,
            x,
            4.5,
            2.0,
            0.36,
            RGBColor(0xEC, 0xEF, 0xF1),
            color=DGRAY,
            size=9,
            bold=True,
        )
        text(s, sub, x, 4.9, 2.0, 0.24, size=8.5, color=MGRAY, align=PP_ALIGN.CENTER)
        x += 2.06
    takeaway(
        s,
        "If a query joins anything other than bn_id_xref or bn_id_persistence, "
        "it probably wants one of those two instead.",
    )
    seealso(s, "Every column:", "identity_hub_data_dictionary.docx", 0.6, 5.72)
    footer(s, n)
    return s


def s_hub_practice(prs, n):
    s = blank(prs)
    bg(s, LGRAY)
    header(
        s,
        "in practice",
        "Asking the hub the right question",
        "Fourteen worked queries. Each states the trap before the SQL.",
    )
    qs = [
        (
            "How large is the audience, honestly",
            "Separates people from browsers. The first query anyone should run.",
        ),
        (
            "The identifier-spread health check",
            "Catches a shared or invalid identifier before it distorts a count.",
        ),
        (
            "What identifiers we hold on known people",
            "What can actually be activated, by channel.",
        ),
    ]
    y = 1.55
    for title, body in qs:
        rect(s, 0.6, y, 8.1, 0.95, WHITE, line=BORDER)
        rect(s, 0.6, y, 0.06, 0.95, NAVY)
        text(s, title, 0.9, y + 0.14, 7.5, 0.3, size=13, bold=True, color=DGRAY)
        text(s, body, 0.9, y + 0.47, 7.5, 0.35, size=11, color=MGRAY)
        y += 1.1
    rect(s, 9.05, 1.55, 3.75, 2.85, WHITE, line=TEAL, line_pt=1.25)
    rect(s, 9.05, 1.55, 3.75, 0.06, TEAL)
    text(
        s, "THE RULE THAT MATTERS", 9.3, 1.75, 3.3, 0.26, size=10, bold=True, color=TEAL
    )
    text(
        s,
        "is_bot = FALSE",
        9.3,
        2.05,
        3.3,
        0.4,
        size=17,
        bold=True,
        color=DGRAY,
        font="Consolas",
    )
    text(
        s,
        "Bot traffic is a real share of what reaches the graph. Any people count "
        "without this filter is inflated, and the inflation is not uniform across "
        "conditions or sites.",
        9.3,
        2.55,
        3.3,
        1.6,
        size=11,
        color=MGRAY,
    )
    text(s, "Full query set:", 0.6, 5.0, 2.2, 0.3, size=12, color=MGRAY)
    text(
        s,
        "identity_hub_bi_queries.docx",
        2.15,
        5.0,
        4.2,
        0.3,
        size=12,
        bold=True,
        color=TEAL,
        font="Consolas",
        link="identity_hub_bi_queries.docx",
    )
    takeaway(
        s,
        "Every query in that document leads with the mistake it prevents. "
        "The rules are the point; the SQL is the easy part.",
    )
    footer(s, n)
    return s


def s_handoff(prs, n):
    s = blank(prs)
    bg(s, LGRAY)
    header(
        s,
        "the join",
        "bn_id is the handover",
        "Everything after this point is keyed on the identity the hub produced.",
    )
    rect(s, 0.7, 2.1, 3.3, 1.7, NAVY)
    text(s, "IDENTITY HUB", 0.95, 2.35, 2.9, 0.3, size=11, bold=True, color=TEAL)
    text(s, "resolves", 0.95, 2.68, 2.9, 0.35, size=16, bold=True, color=WHITE)
    text(
        s,
        "34.2M identifiers\ninto 7.51M people",
        0.95,
        3.05,
        2.9,
        0.6,
        size=11,
        color=SKY,
    )
    arrow(s, 4.25, 2.65, 1.15, 0.6, TEAL)
    chip(s, "bn_id", 4.35, 2.15, 0.95, 0.4, TEAL, size=11)
    rect(s, 5.65, 2.1, 3.3, 1.7, PURPLE)
    text(
        s,
        "PROFILE DATABASE",
        5.9,
        2.35,
        2.9,
        0.3,
        size=11,
        bold=True,
        color=RGBColor(0xE1, 0xBE, 0xE7),
    )
    text(s, "describes", 5.9, 2.68, 2.9, 0.35, size=16, bold=True, color=WHITE)
    text(
        s,
        "each person across\neight source systems",
        5.9,
        3.05,
        2.9,
        0.6,
        size=11,
        color=RGBColor(0xE1, 0xBE, 0xE7),
    )
    arrow(s, 9.2, 2.65, 1.15, 0.6, GREEN)
    rect(s, 10.6, 2.1, 2.2, 1.7, GREEN)
    text(s, "BI", 10.85, 2.4, 1.8, 0.35, size=11, bold=True, color=WHITE)
    text(
        s,
        "audiences,\nreporting,\nactivation",
        10.85,
        2.75,
        1.8,
        0.9,
        size=11.5,
        color=WHITE,
        bold=True,
    )
    warn = (
        "A bn_id can be superseded when two clusters merge. If you cached one "
        "from an earlier extract, resolve it through bn_id_persistence before "
        "joining -- otherwise the row silently disappears."
    )
    rect(s, 0.7, 4.35, 12.1, 0.95, AMBER_BG, line=GOLD)
    text(
        s, "THE ONE TRAP HERE", 0.95, 4.48, 4, 0.26, size=10, bold=True, color=AMBER_TX
    )
    text(s, warn, 0.95, 4.73, 11.6, 0.5, size=12, bold=True, color=AMBER_TX)
    takeaway(
        s,
        "One join key, one direction of travel. The hub never reads from the "
        "profile database.",
    )
    footer(s, n)
    return s


def s_sources(prs, n):
    s = blank(prs)
    bg(s, LGRAY)
    header(
        s,
        "inputs",
        "What the profile database merges",
        "Eight source systems, each contributing a different kind of truth.",
        accent=RGBColor(0xCE, 0x93, 0xD8),
    )
    srcs = [
        ("Mailchimp", "subscription, opens, clicks"),
        ("GA4", "sessions, pages, referral"),
        ("WordPress", "registration, forum, SSO"),
        ("LimeSurvey", "declared condition, demographics"),
        ("NPI registry", "clinician credential, specialty"),
        ("SurveyEngine", "the new registration path"),
        ("AIM", "clickstream engagement"),
        ("Ad platforms", "attribution"),
    ]
    x, y = 0.6, 1.55
    for i, (name, what) in enumerate(srcs):
        if i == 4:
            x, y = 0.6, 2.62
        rect(s, x, y, 2.95, 0.92, WHITE, line=BORDER)
        rect(s, x, y, 2.95, 0.05, PURPLE)
        text(s, name, x + 0.18, y + 0.14, 2.6, 0.28, size=12.5, bold=True, color=DGRAY)
        text(s, what, x + 0.18, y + 0.45, 2.6, 0.4, size=10, color=MGRAY)
        x += 3.06
    arrow(s, 6.15, 3.75, 1.0, 0.5, PURPLE)
    rect(s, 3.9, 4.4, 5.5, 1.0, PURPLE)
    text(
        s,
        "profile_core",
        4.15,
        4.55,
        3.5,
        0.35,
        size=17,
        bold=True,
        color=WHITE,
        font="Consolas",
    )
    text(
        s,
        "one row per person, 131 columns",
        4.15,
        4.93,
        5.0,
        0.3,
        size=11,
        color=RGBColor(0xE1, 0xBE, 0xE7),
    )
    takeaway(
        s,
        "Different sources disagree. Where they do, the database records which "
        "one won and why -- that lineage is what makes a number defensible.",
        colour=PURPLE,
        band=RGBColor(0xF3, 0xE5, 0xF5),
    )
    footer(s, n)
    return s


def s_declared(prs, n):
    s = blank(prs)
    bg(s, LGRAY)
    header(
        s,
        "the distinction that matters most",
        "Declared is not the same as inferred",
        'Both are useful. Only one of them supports the word "patients".',
        accent=RGBColor(0xCE, 0x93, 0xD8),
    )
    rect(s, 0.6, 1.55, 6.0, 3.6, WHITE, line=GREEN, line_pt=1.5)
    rect(s, 0.6, 1.55, 6.0, 0.07, GREEN)
    text(s, "DECLARED", 0.9, 1.78, 4, 0.3, size=12, bold=True, color=GREEN)
    text(s, "They told us", 0.9, 2.08, 4, 0.38, size=19, bold=True, color=DGRAY)
    bullets(
        s,
        [
            "Survey answer, registration form, newsletter chosen",
            "Recorded with the source that supplied it",
            "Supports the words patient, caregiver, diagnosed",
            "Safe to use in a claim to a client",
        ],
        0.9,
        2.6,
        5.4,
        2.2,
        size=12,
    )
    rect(s, 7.0, 1.55, 5.85, 3.6, WHITE, line=GOLD, line_pt=1.5)
    rect(s, 7.0, 1.55, 5.85, 0.07, GOLD)
    text(s, "INFERRED", 7.3, 1.78, 4, 0.3, size=12, bold=True, color=AMBER_TX)
    text(s, "We worked it out", 7.3, 2.08, 5, 0.38, size=19, bold=True, color=DGRAY)
    bullets(
        s,
        [
            "Reading behaviour, site visited, content affinity",
            "Carries a confidence, typically 0.5",
            "Supports interested in, never diagnosed with",
            "Excellent for targeting, wrong for a claim",
        ],
        7.3,
        2.6,
        5.25,
        2.2,
        size=12,
    )
    takeaway(
        s,
        "preferred_condition_source tells you which you are holding. Check it "
        "before any sentence containing the word patients.",
        colour=GOLD,
        band=AMBER_BG,
    )
    footer(s, n)
    return s


def s_funnel(prs, n):
    s = blank(prs)
    bg(s, LGRAY)
    header(
        s,
        "sizing the audience",
        "Each number is smaller than the one above it",
        "Most mistakes come from using a wider number than the question deserves.",
        accent=RGBColor(0xCE, 0x93, 0xD8),
    )
    rows = [
        (
            "7,495,014",
            "profiles in the database",
            "Includes cookie-only records. Browsers, not people.",
            SLATE,
        ),
        (
            "798,780",
            "known people",
            "is_known_person -- someone you could actually name.",
            NAVY,
        ),
        (
            "498,114",
            "with a declared or inferred role",
            "patient, HCP, caregiver, family, other.",
            TEAL,
        ),
        ("407,802", "mailable", "Subscribed, opted in, consent not denied.", GREEN),
    ]
    y = 1.5
    for i, (num, label, note, colour) in enumerate(rows):
        w = 12.2 - i * 1.15
        rect(s, 0.6 + i * 0.57, y, w, 1.02, WHITE, line=BORDER)
        rect(s, 0.6 + i * 0.57, y, 0.08, 1.02, colour)
        text(
            s,
            num,
            0.9 + i * 0.57,
            y + 0.16,
            2.6,
            0.55,
            size=26,
            bold=True,
            color=colour,
        )
        text(
            s,
            label,
            3.55 + i * 0.57,
            y + 0.14,
            4.0,
            0.32,
            size=13.5,
            bold=True,
            color=DGRAY,
        )
        text(s, note, 3.55 + i * 0.57, y + 0.52, 7.4, 0.35, size=11, color=MGRAY)
        y += 1.16
    takeaway(
        s,
        "Quoting 7.50M as an audience overstates it about nine times. "
        "is_known_person is the first filter, not an afterthought.",
        colour=GOLD,
        band=AMBER_BG,
    )
    footer(s, n)
    return s


def s_roles(prs, n):
    s = blank(prs)
    bg(s, LGRAY)
    header(
        s,
        "who they are",
        "Roles and conditions",
        "Roles overlap by design. Conditions are the highest-value field we hold.",
        accent=RGBColor(0xCE, 0x93, 0xD8),
    )
    text(s, "ROLES  (known people)", 0.6, 1.5, 5, 0.26, size=10, bold=True, color=MGRAY)
    roles = [
        ("HCP", 318_164, PURPLE),
        ("Patient", 73_478, TEAL),
        ("Family / friend", 19_098, NAVY),
        ("Other", 11_360, SLATE),
        ("Caregiver", 10_105, GREEN),
    ]
    top = roles[0][1]
    y = 1.85
    for name, cnt, colour in roles:
        rect(s, 2.5, y, 3.6 * (cnt / top), 0.4, colour)
        text(s, name, 0.6, y + 0.05, 1.85, 0.3, size=11.5, bold=True, color=DGRAY)
        text(s, f"{cnt:,}", 6.25, y + 0.06, 1.3, 0.3, size=11, bold=True, color=colour)
        y += 0.55
    text(
        s,
        "Roles are independent flags, not a pie chart. About 3,300 people hold two "
        "or more. Use primary_role when you need mutually exclusive buckets.",
        0.6,
        4.6,
        6.9,
        0.6,
        size=11,
        italic=True,
        color=MGRAY,
    )
    text(
        s,
        "TOP CONDITIONS  (known people)",
        7.7,
        1.5,
        5,
        0.26,
        size=10,
        bold=True,
        color=MGRAY,
    )
    conds = [
        ("Multiple Sclerosis", 59_334),
        ("Parkinsons Disease", 58_507),
        ("ALS", 24_633),
        ("Pulmonary Fibrosis", 20_475),
        ("Ehlers-Danlos Syndrome", 17_046),
    ]
    ctop = conds[0][1]
    y = 1.85
    for name, cnt in conds:
        rect(s, 10.15, y, 2.0 * (cnt / ctop), 0.4, PURPLE)
        text(s, name, 7.7, y + 0.05, 2.4, 0.3, size=11, bold=True, color=DGRAY)
        text(
            s,
            f"{cnt:,}",
            12.25,
            y + 0.06,
            0.9,
            0.3,
            size=10.5,
            bold=True,
            color=PURPLE,
            align=PP_ALIGN.RIGHT,
        )
        y += 0.55
    text(
        s,
        "Check preferred_condition_source before calling any of these patients.",
        7.7,
        4.6,
        5.1,
        0.5,
        size=11,
        italic=True,
        color=MGRAY,
    )
    takeaway(
        s,
        "HCP is the largest role by a wide margin -- and the number most often "
        "quoted wrongly. The next slide is why.",
        colour=PURPLE,
        band=RGBColor(0xF3, 0xE5, 0xF5),
    )
    seealso(s, "Field detail:", "profile_db_table_reference.docx", 0.6, 5.75)
    footer(s, n)
    return s


def s_hcp(prs, n):
    s = blank(prs)
    bg(s, LGRAY)
    header(
        s,
        "the number most often quoted wrongly",
        '"How many HCPs do we have?" has four answers',
        "They differ by 16x. The question decides which one is honest.",
        accent=RED,
    )
    cards = [
        ("380,344", "VERIFIED", "Real, current credential.\nNOT an audience.", SLATE),
        ("102,217", "ENGAGED", "Has ever produced a\npageview, open or click.", TEAL),
        (
            "23,019",
            "REACHABLE",
            "Verified AND mailable.\nUse this commercially.",
            GREEN,
        ),
        ("13,934", "LIVE", "Active on email in the\nlast 90 days.", NAVY),
    ]
    x = 0.6
    for num, label, body, colour in cards:
        rect(s, x, 1.55, 2.95, 2.5, WHITE, line=BORDER)
        rect(s, x, 1.55, 2.95, 0.08, colour)
        text(s, label, x + 0.2, 1.78, 2.5, 0.26, size=10.5, bold=True, color=colour)
        text(s, num, x + 0.2, 2.05, 2.6, 0.6, size=27, bold=True, color=DGRAY)
        text(s, body, x + 0.2, 2.75, 2.6, 1.1, size=11, color=MGRAY)
        x += 3.06
    rect(s, 0.6, 4.25, 12.25, 1.05, REDBG, line=RED)
    text(
        s,
        "WHY THE FIRST NUMBER IS MISLEADING",
        0.85,
        4.38,
        6,
        0.26,
        size=10,
        bold=True,
        color=RED,
    )
    text(
        s,
        "About 73 percent of verified HCPs -- 278,127 people -- came from the "
        "federal NPI registry rather than from our audience. 269,749 of them hold "
        "nothing but an email address and an NPI number: no browser, no session, no "
        "subscription. They are genuine clinicians who have simply never read "
        "anything we publish. Quoting 380,344 as reachable overstates it about "
        "sixteen times.",
        0.85,
        4.63,
        11.7,
        0.6,
        size=11.5,
        bold=True,
        color=DGRAY,
    )
    takeaway(
        s,
        "For anything commercial or external, use 23,019. It is the only figure "
        'that survives "can you actually reach them?"',
        colour=GREEN,
        band=RGBColor(0xE8, 0xF5, 0xE9),
    )
    footer(s, n)
    return s


def s_profile_practice(prs, n):
    s = blank(prs)
    bg(s, LGRAY)
    header(
        s,
        "in practice",
        "Asking the profile database the right question",
        "Fifteen worked queries, each with the rule that stops a number being quoted wrongly.",
        accent=RGBColor(0xCE, 0x93, 0xD8),
    )
    qs = [
        (
            "Known versus unknown, and what we hold on each",
            "The first filter. Everything else is a subset of this.",
        ),
        (
            "Audience by condition, declared and inferred kept apart",
            "Run before quoting any condition audience to anyone.",
        ),
        (
            "The HCP funnel by condition",
            "Verified, engaged, reachable, live -- side by side.",
        ),
        (
            "Growth: real signups only",
            "Registration-grade dates. Excludes cookie sightings.",
        ),
    ]
    y = 1.55
    for title, body in qs:
        rect(s, 0.6, y, 8.1, 0.86, WHITE, line=BORDER)
        rect(s, 0.6, y, 0.06, 0.86, PURPLE)
        text(s, title, 0.9, y + 0.11, 7.6, 0.3, size=12.5, bold=True, color=DGRAY)
        text(s, body, 0.9, y + 0.42, 7.6, 0.32, size=10.5, color=MGRAY)
        y += 0.98
    rect(s, 9.05, 1.55, 3.8, 3.35, WHITE, line=PURPLE, line_pt=1.25)
    rect(s, 9.05, 1.55, 3.8, 0.06, PURPLE)
    text(
        s,
        "WHERE THE FLAGS LIVE",
        9.3,
        1.75,
        3.3,
        0.26,
        size=10,
        bold=True,
        color=PURPLE,
    )
    text(
        s,
        "profile_metrics",
        9.3,
        2.05,
        3.4,
        0.4,
        size=15,
        bold=True,
        color=DGRAY,
        font="Consolas",
    )
    text(
        s,
        "One view holding every audience flag as a clean boolean: "
        "is_known_person, is_verified_hcp, is_engaged_hcp, is_mailable, "
        "is_active_email_90d.\n\nBuilt so nobody has to re-derive a definition "
        "and get it slightly different.",
        9.3,
        2.5,
        3.35,
        2.2,
        size=11,
        color=MGRAY,
    )
    text(s, "Full query set:", 0.6, 5.5, 2.2, 0.3, size=12, color=MGRAY)
    text(
        s,
        "profile_db_bi_queries.docx",
        2.15,
        5.5,
        4.2,
        0.3,
        size=12,
        bold=True,
        color=TEAL,
        font="Consolas",
        link="profile_db_bi_queries.docx",
    )
    takeaway(
        s,
        "If a definition is worth arguing about, it belongs in profile_metrics, "
        "not in each analyst's query.",
        colour=PURPLE,
        band=RGBColor(0xF3, 0xE5, 0xF5),
    )
    seealso(s, "View columns:", "profile_db_data_dictionary.docx", 0.6, 5.85, 4.6)
    footer(s, n)
    return s


def s_recent(prs, n):
    s = blank(prs)
    bg(s, LGRAY)
    header(
        s,
        "read before quoting",
        "What changed recently",
        "Four corrections landed in August. Each prevents a specific wrong number.",
        accent=RED,
    )
    items = [
        (
            "The HCP number has three sizes",
            "380,344 verified, 102,217 engaged, 23,019 reachable. Use the last one commercially.",
        ),
        (
            "is_engaged_hcp says ever, not lately",
            "No time bound. Only 47,802 were active in the last 90 days. It is a floor, not a measurement.",
        ),
        (
            "Mailable counts rose, nobody was newly opted in",
            "A cookie-consent signal was gating the email metrics. 5,220 opted-in people were wrongly excluded.",
        ),
        (
            "A lapsed credential is not a lapsed person",
            "4,170 retired NPIs left is_verified_hcp but keep their history. 95 percent still read.",
        ),
    ]
    y = 1.55
    for title, body in items:
        rect(s, 0.6, y, 12.25, 1.05, WHITE, line=BORDER)
        rect(s, 0.6, y, 0.07, 1.05, RED)
        text(s, title, 0.95, y + 0.14, 5.3, 0.32, size=13, bold=True, color=RED)
        text(s, body, 6.4, y + 0.22, 6.2, 0.7, size=11.5, color=DGRAY)
        y += 1.17
    text(
        s,
        "All four are marked in red inside",
        0.6,
        6.3,
        3.6,
        0.3,
        size=11.5,
        color=MGRAY,
    )
    text(
        s,
        "profile_db_bi_queries.docx",
        3.95,
        6.3,
        4.0,
        0.3,
        size=11.5,
        bold=True,
        color=TEAL,
        font="Consolas",
        link="profile_db_bi_queries.docx",
    )
    footer(s, n)
    return s


def s_trust(prs, n):
    s = blank(prs)
    bg(s, LGRAY)
    header(
        s,
        "confidence",
        "Why you can trust these numbers",
        "The platform rebuilds daily and refuses to publish a build it cannot vouch for.",
        accent=GREEN,
    )
    stats = [
        ("16", "acceptance gates", "run on every single build", GREEN),
        (
            "0",
            "published on hard failure",
            "a failed gate blocks the release outright",
            RED,
        ),
        (
            "13",
            "months of snapshots",
            "monthly forensic copies, so a late question is answerable",
            TEAL,
        ),
        (
            "Daily",
            "identity + profile refresh",
            "both systems rebuild every morning",
            NAVY,
        ),
    ]
    x = 0.6
    for big, label, sub, colour in stats:
        rect(s, x, 1.6, 2.95, 2.15, WHITE, line=BORDER)
        rect(s, x, 1.6, 2.95, 0.08, colour)
        text(s, big, x + 0.22, 1.85, 2.6, 0.7, size=32, bold=True, color=colour)
        text(s, label, x + 0.22, 2.6, 2.6, 0.3, size=12, bold=True, color=DGRAY)
        text(s, sub, x + 0.22, 2.92, 2.6, 0.75, size=10.5, color=MGRAY)
        x += 3.06
    rect(s, 0.6, 4.1, 12.25, 1.35, WHITE, line=GREEN, line_pt=1.25)
    text(
        s,
        "WHAT A GATE ACTUALLY DOES",
        0.9,
        4.25,
        5,
        0.26,
        size=10,
        bold=True,
        color=GREEN,
    )
    text(
        s,
        "A build that loses too many people, drops a fill rate, or receives a "
        "truncated identity graph is stopped before it reaches you. Production "
        "keeps serving the previous good release while somebody investigates. "
        "The numbers on a dashboard are never a half-finished build.",
        0.9,
        4.55,
        11.6,
        0.8,
        size=12.5,
        color=DGRAY,
    )
    takeaway(
        s,
        "If a figure looks wrong, it is far more likely a definition question "
        "than a broken build. Start with the BI query documents.",
        colour=GREEN,
        band=RGBColor(0xE8, 0xF5, 0xE9),
    )
    footer(s, n)
    return s


def s_lookup(prs, n):
    s = blank(prs)
    bg(s, LGRAY)
    header(
        s,
        "reference",
        "The eight documents",
        "Read in this order. Click any name to open it.",
        accent=TEAL,
    )
    y = 1.5
    for fname, title, colour, desc in DOCS:
        rect(s, 0.6, y, 12.25, 0.63, WHITE, line=BORDER)
        rect(s, 0.6, y, 0.07, 0.63, colour)
        text(s, title, 0.9, y + 0.06, 4.05, 0.3, size=12, bold=True, color=colour)
        text(s, desc, 4.95, y + 0.08, 5.55, 0.45, size=10, color=MGRAY)
        text(
            s,
            fname,
            10.6,
            y + 0.09,
            2.2,
            0.3,
            size=9,
            color=TEAL,
            font="Consolas",
            link=fname,
            align=PP_ALIGN.RIGHT,
        )
        y += 0.71
    text(
        s,
        "Links are relative -- keep this deck in the same folder as the documents. "
        f"All figures verified {AS_OF}.",
        0.6,
        7.05,
        10,
        0.3,
        size=9,
        italic=True,
        color=MGRAY,
    )
    text(s, str(n), 12.55, 7.05, 0.35, 0.28, size=9, color=MGRAY, align=PP_ALIGN.RIGHT)
    return s


def main() -> int:
    EXPORTS.mkdir(parents=True, exist_ok=True)
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    s_title(prs)
    builders = [
        s_two_halves,
        s_problem,
        s_identifiers,
        s_matching,
        s_hub_output,
        s_hub_practice,
        s_handoff,
        s_sources,
        s_declared,
        s_funnel,
        s_roles,
        s_hcp,
        s_profile_practice,
        s_recent,
        s_trust,
        s_lookup,
    ]
    n = 2
    for b in builders:
        b(prs, n)
        n += 1
    prs.save(str(OUT))
    print(f"[OK] Wrote {OUT.relative_to(REPO)}  ({len(prs.slides._sldIdLst)} slides)")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
